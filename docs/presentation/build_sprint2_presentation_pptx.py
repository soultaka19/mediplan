# -*- coding: utf-8 -*-
"""Genere le support de PRESENTATION (revue) du Sprint 2 de MediPlan (PowerPoint).

Consigne : presentation du 29-30 juillet 2026 portant sur CE QUI A ETE REALISE
au Sprint 2 (et non le plan). Structure calquee sur les elements demandes par la
professeure : rappel du projet, objectif du sprint, developpement de la solution,
demonstration, tests et validation, difficultes, avancement Jira, prochaines etapes.

Voir Sprint2-presentation.md pour le decoupage, le minutage et le script de demo.

Usage : python docs/presentation/build_sprint2_presentation_pptx.py
Sortie : docs/presentation/MediPlan-Sprint2-Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette (identique aux supports precedents, continuite visuelle) ---
BLEU = RGBColor(0x15, 0x4D, 0x8C)
BLEU_CLAIR = RGBColor(0x1E, 0x88, 0xE5)
SARCELLE = RGBColor(0x00, 0x89, 0x7B)
GRIS = RGBColor(0x40, 0x40, 0x40)
GRIS_CLAIR = RGBColor(0x6B, 0x6B, 0x6B)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
FOND = RGBColor(0xF4, 0xF7, 0xFB)
BLEU_PALE = RGBColor(0xCF, 0xE3, 0xFF)
BLEU_PALE2 = RGBColor(0x9F, 0xC8, 0xF5)
AMBRE = RGBColor(0xB5, 0x6A, 0x00)
VERT = RGBColor(0x2E, 0x7D, 0x32)

LARGEUR = Inches(13.333)
HAUTEUR = Inches(7.5)

prs = Presentation()
prs.slide_width = LARGEUR
prs.slide_height = HAUTEUR
BLANK = prs.slide_layouts[6]


def fond(slide, couleur=BLANC):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = couleur


def rect(slide, x, y, w, h, couleur):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = couleur
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def zone_texte(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def para(tf, texte, taille=18, couleur=GRIS, gras=False, puce=False,
         espace_avant=6, niveau=0, aligne=PP_ALIGN.LEFT, premier=False):
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    p.alignment = aligne
    p.level = niveau
    p.space_after = Pt(espace_avant)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = (("•  " if puce else "") + texte)
    run.font.size = Pt(taille)
    run.font.bold = gras
    run.font.color.rgb = couleur
    run.font.name = "Calibri"
    return p


def notes(slide, texte):
    slide.notes_slide.notes_text_frame.text = texte


def bandeau_titre(slide, numero, titre, sous_titre=None):
    rect(slide, 0, 0, LARGEUR, Inches(1.25), BLEU)
    rect(slide, 0, Inches(1.25), LARGEUR, Pt(6), SARCELLE)
    past = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.32),
                                  Inches(0.62), Inches(0.62))
    past.fill.solid(); past.fill.fore_color.rgb = SARCELLE
    past.line.fill.background(); past.shadow.inherit = False
    pf = past.text_frame; pf.word_wrap = False
    pp = pf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); r.text = str(numero); r.font.size = Pt(22)
    r.font.bold = True; r.font.color.rgb = BLANC
    tb, tf = zone_texte(slide, Inches(1.25), Inches(0.18), Inches(11.6), Inches(1.0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, titre, taille=27, couleur=BLANC, gras=True, premier=True, espace_avant=0)
    if sous_titre:
        para(tf, sous_titre, taille=14, couleur=BLEU_PALE, espace_avant=0)


def pied(slide, orateur, duree):
    tb, tf = zone_texte(slide, Inches(0.45), Inches(7.0), Inches(12.4), Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = f"Orateur : {orateur}"
    r.font.size = Pt(11); r.font.color.rgb = GRIS_CLAIR; r.font.italic = True
    tb2, tf2 = zone_texte(slide, Inches(10.0), Inches(7.0), Inches(2.9), Inches(0.4))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"⏱  {duree}"
    r2.font.size = Pt(11); r2.font.color.rgb = SARCELLE; r2.font.bold = True


def carte(slide, x, y, w, h, titre, lignes, couleur_titre=BLEU, taille_ligne=13):
    rect(slide, x, y, w, h, FOND)
    rect(slide, x, y, Pt(6), h, couleur_titre)
    tb, tf = zone_texte(slide, x + Inches(0.25), y + Inches(0.12),
                        w - Inches(0.4), h - Inches(0.2))
    para(tf, titre, taille=16, couleur=couleur_titre, gras=True, premier=True,
         espace_avant=4)
    for l in lignes:
        para(tf, l, taille=taille_ligne, couleur=GRIS, puce=bool(l.strip()), espace_avant=3)


def chiffre(slide, x, y, largeur, valeur, libelle, couleur):
    rect(slide, x, y, largeur, Inches(1.5), couleur)
    tb, tf = zone_texte(slide, x, y + Inches(0.08), largeur, Inches(1.4))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, valeur, taille=38, couleur=BLANC, gras=True, premier=True,
         aligne=PP_ALIGN.CENTER, espace_avant=0)
    para(tf, libelle, taille=12, couleur=BLANC, aligne=PP_ALIGN.CENTER, espace_avant=0)


# =========================================================================
# DIAPO 1 — Titre
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(3.05), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(1.3), Inches(11.3), Inches(1.4))
para(tf, "MediPlan", taille=58, couleur=BLANC, gras=True, premier=True,
     aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.7))
para(tf, "Plateforme web de gestion des rendez-vous médicaux", taille=22,
     couleur=BLEU_PALE, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(3.25), Inches(11.3), Inches(1.4))
para(tf, "Présentation du Sprint 2", taille=28, couleur=BLANC,
     gras=True, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=8)
para(tf, "Revue de réalisation — ce qui a été développé, testé et intégré", taille=18,
     couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.5))
para(tf, "Équipe : Souleymane DIALLO  ·  Zakaria Lahouiri  ·  Larbi Saib", taille=18,
     couleur=BLANC, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=4)
para(tf, "Projet intégrateur — Programmation informatique  ·  Collège La Cité  ·  Printemps 2026",
     taille=14, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=2)
para(tf, "Présentation des 29–30 juillet 2026", taille=13,
     couleur=BLEU_PALE, aligne=PP_ALIGN.CENTER, espace_avant=6)
notes(s, "SOULEYMANE — 20 s\n\n"
         "ACCROCHE (à dire telle quelle) :\n"
         "« Dans une clinique, quand deux personnes décrochent le téléphone en même temps, "
         "il arrive que le même créneau soit vendu deux fois. C'est un agenda papier, une "
         "rature, un patient qui se déplace pour rien. »\n\n"
         "[ MARQUER UN TEMPS — 1 seconde ]\n\n"
         "« Le Sprint 2, c'est le sprint où ça devient impossible dans MediPlan. Bonjour, "
         "nous sommes l'équipe MediPlan : Souleymane, Zakaria, Larbi. En dix minutes, nous "
         "allons vous montrer ce que nous avons développé, testé et intégré. »\n\n"
         "→ TRANSITION : « D'abord, un rappel rapide du projet. »\n\n"
         "RAPPEL : on ouvre sur le problème métier, pas sur la technologie.")

# =========================================================================
# DIAPO 2 — Rappel du projet
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 1, "Rappel du projet", "MediPlan — gérer les rendez-vous d'une clinique")
carte(s, Inches(0.45), Inches(1.65), Inches(6.1), Inches(2.5),
      "Problématique",
      ["Dans une clinique, la prise de rendez-vous par téléphone est",
       "   manuelle : agendas papier, doubles réservations, oublis",
       "Pas de vue claire de l'activité du jour pour la réception",
       "Données de patients sensibles à protéger et à cloisonner"], AMBRE)
carte(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(2.5),
      "Solution proposée",
      ["Application web : la réception gère les rendez-vous à la place",
       "   des patients (qui appellent au téléphone)",
       "Disponibilités des médecins → créneaux réservables générés",
       "Flux du jour en temps réel + contrôle d'accès par rôle (RBAC)"], SARCELLE)
carte(s, Inches(0.45), Inches(4.4), Inches(12.4), Inches(2.35),
      "Utilisateurs / clients visés",
      ["Réception / administrateur de clinique — utilisateur principal : crée et suit les rendez-vous",
       "Médecin — consulte ses disponibilités et le flux de sa journée",
       "Patient « léger » — créé au comptoir, sans compte ni mot de passe (il téléphone, il ne s'inscrit pas)",
       "Super-administrateur — supervise l'ensemble des cliniques",
       " ",
       "Stack : Angular (frontend) · NestJS (backend) · PostgreSQL (base) · Docker (environnement)"], BLEU)
pied(s, "Souleymane", "1 min 30")
notes(s, "SOULEYMANE — 1 min 30\n\n"
         "« Le problème : la prise de rendez-vous se fait au téléphone, à la main. Agendas "
         "papier, doubles réservations, oublis. Et la réception n'a aucune vue claire de sa "
         "journée. »\n\n"
         "« Notre solution, c'est une application web où la réception gère les rendez-vous à "
         "la place des patients. Et c'est notre choix le plus structurant : notre "
         "utilisateur principal n'est pas le patient — c'est la réception. Le patient "
         "appelle, il ne s'inscrit pas en ligne. »\n\n"
         "« Quatre profils : la réception, qui fait tout ; le médecin, qui consulte sa "
         "journée ; le patient dit léger, créé au comptoir sans compte ni mot de passe ; et "
         "un super-administrateur. Le tout en Angular, NestJS, PostgreSQL, dans Docker. »\n\n"
         "ATTENTION : ne pas détailler la stack — Zakaria y revient à la diapo 4.\n\n"
         "→ TRANSITION : « Voilà le projet. Maintenant, qu'est-ce qu'on s'était engagé à "
         "livrer sur ce sprint ? »")

# =========================================================================
# DIAPO 3 — Objectif du Sprint 2
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 2, "Objectif du Sprint 2", "Ce que l'équipe voulait accomplir")
rect(s, Inches(0.45), Inches(1.6), Inches(12.4), Inches(1.05), SARCELLE)
tb, tf = zone_texte(s, Inches(0.75), Inches(1.7), Inches(11.8), Inches(0.9))
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "Objectif", taille=12, couleur=BLEU_PALE, gras=True, premier=True, espace_avant=2)
para(tf, "Une réception se connecte de façon sécurisée, définit les disponibilités des médecins "
         "et réserve un rendez-vous pour un patient — sans jamais créer de double réservation.",
     taille=16, couleur=BLANC, gras=True, espace_avant=0)
carte(s, Inches(0.45), Inches(2.9), Inches(6.1), Inches(3.85),
      "Fonctionnalités visées",
      ["Authentification complète + contrôle d'accès (RBAC 4 rôles)",
       "Modéliser le « patient léger » géré par la réception",
       "Définir les disponibilités des médecins",
       "Générer automatiquement les créneaux réservables",
       "Réserver un rendez-vous (réception) sans double-réservation",
       "Suivre le flux clinique du jour (statuts de consultation)",
       "Interface soignée + mode sombre (design system)"], BLEU)
carte(s, Inches(6.75), Inches(2.9), Inches(6.1), Inches(3.85),
      "Pourquoi ce périmètre",
      ["C'est la tranche verticale la plus utile : de la connexion",
       "   jusqu'à un rendez-vous réel dans l'agenda",
       "Elle couvre le cœur métier (anti-double-réservation) et la",
       "   sécurité (données de santé, RBAC)",
       "Chaque fonctionnalité est démontrable de bout en bout",
       " ",
       "Note : notre découpage interne comptait plusieurs sprints ;",
       "   ils ont été alignés sur le calendrier du cours."], SARCELLE)
pied(s, "Souleymane", "1 min")
notes(s, "SOULEYMANE — 1 min — C'EST ICI QU'ON POSE LA PHRASE DU JOUR. DIRE LENTEMENT.\n\n"
         "« Notre objectif tenait en une phrase : une réception se connecte de façon "
         "sécurisée, définit les disponibilités des médecins, et réserve un rendez-vous pour "
         "un patient — sans jamais créer de double réservation. »\n\n"
         "[ TEMPS ]\n\n"
         "« Au Sprint 1, on savait ouvrir une session. Au Sprint 2, la réception fait "
         "tourner une vraie journée de clinique : elle ouvre l'agenda, elle réserve, elle "
         "suit le flux, elle annule. De bout en bout. »\n\n"
         "« Pourquoi ce périmètre-là ? Parce que c'est la tranche verticale la plus utile : "
         "de la connexion jusqu'à un vrai rendez-vous dans l'agenda. Nous avons "
         "volontairement écarté le libre-service patient. Un sprint qui livre un parcours "
         "complet vaut mieux que trois demi-fonctionnalités. »\n\n"
         "→ TRANSITION : « Zakaria va vous expliquer comment tout ça est construit. » "
         "— ET SE TOURNER VERS LUI.")

# =========================================================================
# DIAPO 4 — Developpement : architecture & lien entre les parties
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 3, "Développement — architecture générale",
              "Comment les parties fonctionnent ensemble")

# Trois blocs + fleches (frontend -> backend -> base)
def bloc(x, titre, lignes, couleur):
    rect(s, x, Inches(1.85), Inches(3.5), Inches(2.5), couleur)
    tb, tf = zone_texte(s, x + Inches(0.2), Inches(2.0), Inches(3.15), Inches(2.25))
    para(tf, titre, taille=17, couleur=BLANC, gras=True, premier=True, espace_avant=4,
         aligne=PP_ALIGN.CENTER)
    for l in lignes:
        para(tf, l, taille=12, couleur=BLANC, espace_avant=3, aligne=PP_ALIGN.CENTER)

bloc(Inches(0.55), "Frontend — Angular",
     ["Écrans réception & médecin", "Signals · Material 3", "Design system (tokens)"], BLEU_CLAIR)
bloc(Inches(4.9), "Backend — NestJS",
     ["API REST /api/v1", "JWT + guards RBAC", "Validation par DTO"], SARCELLE)
bloc(Inches(9.25), "Base — PostgreSQL",
     ["TypeORM + migrations", "Index anti-double", "Données cloisonnées"], BLEU)
# fleches
for fx in (Inches(4.35), Inches(8.7)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, fx, Inches(2.9), Inches(0.5), Inches(0.45))
    a.fill.solid(); a.fill.fore_color.rgb = GRIS_CLAIR; a.line.fill.background(); a.shadow.inherit = False

carte(s, Inches(0.45), Inches(4.7), Inches(12.4), Inches(2.05),
      "Le lien entre les parties (un flux de réservation)",
      ["1.  Le navigateur envoie une requête HTTPS au backend, avec le jeton JWT dans l'en-tête (proxy /api → NestJS)",
       "2.  NestJS valide le jeton, vérifie le rôle (guard RBAC) et les données reçues (DTO), puis applique la règle métier",
       "3.  TypeORM écrit en base dans une transaction ; l'index unique partiel garantit qu'un créneau n'est pris qu'une fois",
       "4.  La réponse remonte au frontend, qui met à jour l'affichage via les signals — monorepo unique (frontend + backend)"], BLEU)
pied(s, "Zakaria", "1 min 30")
notes(s, "ZAKARIA — 1 min 30\n\n"
         "« Trois parties dans un seul dépôt : un frontend Angular, un backend NestJS, une "
         "base PostgreSQL. »\n\n"
         "« Suivons une réservation, de bout en bout. Un : le navigateur envoie la requête "
         "au backend, avec le jeton d'authentification dans l'en-tête. Deux : NestJS vérifie "
         "le jeton, vérifie le rôle de l'utilisateur, valide les données reçues, puis "
         "applique la règle métier. Trois : on écrit en base dans une transaction — et c'est "
         "là que se joue le cœur du sprint : un index unique garantit qu'un créneau ne peut "
         "être pris qu'une seule fois. Quatre : la réponse remonte, et l'écran se met à jour "
         "tout seul. »\n\n"
         "LE GESTE : suivre les trois blocs avec la main, de gauche à droite. On doit VOIR "
         "le flux.\n\n"
         "→ TRANSITION : « Regardons maintenant chaque couche de plus près. »")

# =========================================================================
# DIAPO 5 — Developpement : les trois couches en detail
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 4, "Développement — base, backend, frontend",
              "Les choix techniques, couche par couche")
carte(s, Inches(0.45), Inches(1.65), Inches(4.05), Inches(5.0),
      "Base de données",
      ["PostgreSQL, schéma issu de l'ERD",
       "Migrations TypeORM versionnées",
       "Entités : Clinique, Utilisateur,",
       "   Disponibilité, Créneau, RDV",
       "Patient léger : rôle patient,",
       "   mot de passe NULL",
       "Cœur métier :",
       "   index unique partiel sur le",
       "   créneau (hors annulés) →",
       "   pas de double réservation",
       "Cloisonnement par clinic_id"], BLEU)
carte(s, Inches(4.7), Inches(1.65), Inches(4.05), Inches(5.0),
      "Backend — NestJS",
      ["API REST sous /api/v1",
       "Modules : auth, user, clinic,",
       "   availability, appointment",
       "JWT (HS256, 60 min), bcrypt",
       "Guards + @Roles → RBAC",
       "Validation des entrées par DTO",
       "   (class-validator)",
       "Verrouillage de compte (5/15 min)",
       "Transactions + verrous pour la",
       "   réservation concurrente"], SARCELLE)
carte(s, Inches(8.95), Inches(1.65), Inches(3.9), Inches(5.0),
      "Frontend — Angular",
      ["Composants standalone",
       "État réactif par signals",
       "Angular Material 3",
       "Design system : une seule",
       "   source de tokens (couleurs,",
       "   typographie) → cohérence",
       "Intercepteur : jeton ajouté",
       "   automatiquement, gestion 401",
       "Masquage des vues selon le rôle",
       "Mode clair / sombre complet"], BLEU_CLAIR)
pied(s, "Zakaria", "1 min 30")
notes(s, "ZAKARIA — 1 min 30 — NE PAS LIRE LES PUCES. TROIS PHRASES PAR CARTE.\n\n"
         "BASE : « Un schéma issu directement de notre modèle de conception, avec des "
         "migrations versionnées — donc reproductible sur une machine neuve. La pièce "
         "maîtresse, c'est un index unique partiel sur le créneau : les rendez-vous annulés "
         "en sont exclus, donc un créneau annulé redevient réservable. Et le patient léger "
         "est un utilisateur sans mot de passe : il ne peut pas se connecter, par "
         "construction. »\n\n"
         "BACKEND : « Une API REST en modules — authentification, utilisateurs, "
         "disponibilités, rendez-vous. Jeton JWT, mots de passe hachés avec bcrypt, "
         "verrouillage du compte après cinq échecs, et un contrôle d'accès par rôle sur "
         "chaque route. Toutes les entrées sont validées avant d'atteindre la logique "
         "métier. »\n\n"
         "FRONTEND : « Des composants autonomes, un état réactif, et surtout un design "
         "system unique — une seule source de couleurs et de typographie. C'est ce qui fait "
         "qu'on a un mode sombre complet sans avoir retouché un seul écran à la main. »\n\n"
         "→ TRANSITION : « Assez de théorie — Larbi va vous le montrer en vrai. »")

# =========================================================================
# DIAPO 6 — Demonstration
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 5, "Démonstration — ce qui fonctionne", "Parcours réel de la réception")
carte(s, Inches(0.45), Inches(1.65), Inches(7.4), Inches(5.05),
      "Le scénario que nous démontrons en direct",
      ["1.  Connexion réception → l'en-tête affiche « Alice Tremblay »",
       "     (le profil vient du serveur, pas du jeton)",
       "2.  Tableau de bord → le compteur « RDV du jour » est un chiffre réel",
       "3.  Disponibilités → saisir une plage datée + une durée",
       "     → le système génère seul les créneaux réservables",
       "4.  Réserver → choisir médecin, créneau libre, patient, motif",
       "     → le rendez-vous entre dans l'agenda",
       "5.  Flux du jour → cycle : Arrivé → En consultation → Terminé",
       "6.  Annuler un rendez-vous avec motif → le créneau se libère",
       "     et redevient réservable  (fonctionnalité déjà livrée)",
       "7.  RBAC visible : « Utilisateurs » disparaît pour un médecin",
       "8.  Mode sombre en un clic"], SARCELLE, taille_ligne=13)
carte(s, Inches(8.05), Inches(1.65), Inches(4.8), Inches(5.05),
      "Logique de fonctionnement",
      ["Le patient n'agit pas en ligne :",
       "   la réception fait tout pour lui",
       " ",
       "Une disponibilité datée →",
       "   des créneaux → un rendez-vous",
       " ",
       "Un créneau ne peut être pris",
       "   qu'une fois (garanti en base)",
       " ",
       "Chaque écran est protégé selon",
       "   le rôle de l'utilisateur connecté",
       " ",
       "Données de démo réalistes en place",
       "   (2 médecins, 10 patients, RDV",
       "   du jour à statuts variés)"], BLEU)
pied(s, "Larbi (démonstration)", "2 min")
notes(s, "LARBI — 2 min — RÈGLE D'OR : PARLER PENDANT QU'ON CLIQUE, JAMAIS APRÈS.\n\n"
         "1. Connexion réception : « Je me connecte comme la réception. Regardez l'en-tête : "
         "Alice Tremblay, son nom — pas son e-mail. Le profil vient du serveur. »\n\n"
         "2. Tableau de bord : « Le compteur rendez-vous du jour est un vrai chiffre, "
         "calculé en base. »\n\n"
         "3. Disponibilités, créer une plage, puis Voir les créneaux : « Je saisis une plage "
         "datée et une durée de créneau. Le système génère seul les créneaux réservables. Je "
         "ne les ai pas créés un par un. »\n\n"
         "4. RÉSERVER — LE SOMMET, RALENTIR ICI : « Le patient appelle. Je le crée au "
         "comptoir — pas de compte, pas de mot de passe. Je choisis un créneau libre, un "
         "motif, je réserve… » [TEMPS] « …et il est immédiatement dans la journée. C'est le "
         "geste réel d'une réceptionniste. »\n\n"
         "5. Arrivé, Consultation, Terminé : « Et voilà le cycle d'une consultation, tel "
         "qu'il se vit au comptoir. »\n\n"
         "6. Annuler avec motif : « J'annule, avec un motif obligatoire. Le créneau se "
         "libère et redevient réservable. »\n\n"
         "7. Déconnexion puis connexion médecin : « Même application, autre rôle. Regardez "
         "le menu : Utilisateurs a disparu. Le contrôle d'accès se voit à l'écran. »\n\n"
         "8. Mode sombre : « Et le thème sombre, en un clic, sur toute l'application. »\n\n"
         "PRATIQUE : naviguer par le menu de gauche, éviter les rechargements de page.\n"
         "SI ÇA PLANTE : « Le temps que ça revienne, je vous explique ce qui se passe "
         "derrière. » — ne jamais déboguer en direct.\n\n"
         "→ TRANSITION : « Ça, c'est ce que vous voyez. Voyons ce qu'on a vérifié. »")

# =========================================================================
# DIAPO 7 — Tests et validation
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 6, "Tests et validation", "Ce que nous avons vérifié")
chiffre(s, Inches(0.45), Inches(1.6), Inches(2.85), "171", "Tests automatisés verts", VERT)
chiffre(s, Inches(3.55), Inches(1.6), Inches(2.85), "48", "Tests backend (NestJS)", SARCELLE)
chiffre(s, Inches(6.65), Inches(1.6), Inches(2.85), "123", "Tests frontend (Angular)", BLEU_CLAIR)
chiffre(s, Inches(9.75), Inches(1.6), Inches(2.85), "0", "Test en échec", BLEU)
carte(s, Inches(0.45), Inches(3.35), Inches(6.1), Inches(3.4),
      "Tests réalisés",
      ["Backend : services et contrôleurs (auth, disponibilités,",
       "   réservation, flux) + validation des DTO",
       "Anti-double-réservation testé en situation de concurrence",
       "Frontend : composants, services HTTP, guards de rôle,",
       "   intercepteurs, façade d'authentification",
       "Vérification manuelle bout-en-bout dans le navigateur",
       "Migrations rejouées sur une base neuve (reproductible)"], VERT)
carte(s, Inches(6.75), Inches(3.35), Inches(6.1), Inches(3.4),
      "Bogues identifiés → corrigés",
      ["Erreur 500 à la réservation (verrou SQL mal placé sur une",
       "   jointure) → verrou ciblé sur la ligne du RDV",
       "Sélecteur « Médecin » qui ne s'ouvrait pas (l'étiquette",
       "   captait le clic) → corrigé",
       "Rendez-vous du jour vides (fuseau horaire) → dates relatives",
       "En-tête affichant l'e-mail au lieu du nom → repli sur le nom",
       "Motif d'annulation vide accepté → validation renforcée"], AMBRE)
pied(s, "Larbi", "1 min 30")
notes(s, "LARBI — 1 min 30\n\n"
         "« 171 tests automatisés, tous verts : 48 sur le backend, 123 sur le frontend, zéro "
         "en échec. »\n\n"
         "« Le test dont nous sommes le plus fiers : la réservation en concurrence. Deux "
         "réservations lancées sur le même créneau au même instant — une passe, l'autre est "
         "refusée. »\n\n"
         "« Et nous avons trouvé de vrais bogues, que nous avons corrigés : une erreur 500 "
         "causée par un verrou SQL mal placé ; un sélecteur de médecin qui ne s'ouvrait pas "
         "parce que l'étiquette captait le clic ; des rendez-vous du jour vides à cause d'un "
         "problème de fuseau horaire. Tous corrigés. »\n\n"
         "POURQUOI CITER LES BOGUES : une équipe qui nomme ses bogues prouve qu'elle a "
         "réellement testé. Une équipe qui n'en cite aucun n'a testé que le chemin "
         "heureux.\n\n"
         "→ TRANSITION : « Souleymane va vous parler de ce qui a été le plus dur. »")

# =========================================================================
# DIAPO 8 — Difficultes rencontrees et solutions
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 7, "Difficultés rencontrées et solutions",
              "Techniques et organisationnelles")
carte(s, Inches(0.45), Inches(1.65), Inches(6.1), Inches(2.55),
      "⚠  Organisation — dérive d'intégration (vécue)",
      ["Plusieurs branches en parallèle jamais fusionnées",
       "→ deux migrations sur le même horodatage, un module",
       "   de rendez-vous recréé en double",
       "Solution : réintégration manuelle, puis 6 règles (socle",
       "   d'abord, PR sous 72 h, un module par ticket, gates de",
       "   fusion, point à mi-sprint, « terminé » = fusionné dans dev)"], AMBRE)
carte(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(2.55),
      "⚠  Technique — la réservation concurrente",
      ["Comment empêcher deux réceptions de réserver le même",
       "   créneau au même instant ?",
       "Solution : index unique partiel en base + transaction",
       "   avec verrou → c'est PostgreSQL qui tranche la course,",
       "   pas le code applicatif",
       "Vérifié par un test de concurrence"], BLEU_CLAIR)
carte(s, Inches(0.45), Inches(4.45), Inches(12.4), Inches(2.3),
      "Autres difficultés techniques levées",
      ["Verrou SQL « FOR UPDATE » posé sur une jointure externe → erreur 500 : on verrouille la seule ligne du RDV",
       "Dépendance circulaire au démarrage d'Angular (intercepteur ↔ authentification) → initialisation différée",
       "Identifiants de démo invalides (UUID) refusés par la validation → génération d'UUID conformes",
       "Fuseau horaire : rendez-vous « du jour » calculés en référence à Toronto pour rester cohérents"], SARCELLE)
pied(s, "Souleymane", "1 min 30")
notes(s, "SOULEYMANE — 1 min 30 — VOTRE MOMENT. RALENTIR, REGARDER LA SALLE, NE PAS "
         "S'EXCUSER.\n\n"
         "« Notre difficulté la plus sérieuse n'a pas été technique. Elle a été "
         "d'organisation. »\n\n"
         "[ TEMPS ]\n\n"
         "« Au début du sprint, nous avons travaillé chacun sur notre branche, et nous ne "
         "les avons pas fusionnées assez tôt. Résultat : deux migrations sur le même "
         "horodatage, et un module de rendez-vous écrit deux fois par deux personnes "
         "différentes, sans qu'on le sache. Nous avons perdu une journée entière à "
         "réintégrer au lieu de développer. »\n\n"
         "« Nous en avons tiré six règles, que nous appliquons depuis : le code qui touche "
         "la base passe en premier ; aucune branche ne vit plus de trois jours sans être "
         "proposée à la fusion ; un seul module par ticket ; on ne fusionne que si les tests "
         "passent ; un point d'intégration à mi-sprint. Et surtout, notre définition de "
         "terminé a changé : terminé, ça ne veut plus dire codé sur ma machine, ça veut dire "
         "fusionné dans la branche commune. »\n\n"
         "« Côté technique, la vraie question était : comment empêcher deux réceptions de "
         "réserver le même créneau au même instant ? Notre réponse, c'est de ne pas laisser "
         "notre code arbitrer. C'est PostgreSQL qui tranche la course, avec un index unique "
         "et une transaction. Vérifié par un test de concurrence. »\n\n"
         "→ TRANSITION : « Voyons maintenant où ça nous mène dans Jira. »")

# =========================================================================
# DIAPO 9 — Avancement dans Jira
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 8, "Avancement dans Jira", "Répartition et progression")
chiffre(s, Inches(0.45), Inches(1.6), Inches(3.9), "18 / 21", "Tickets Sprint 2 terminés", VERT)
chiffre(s, Inches(4.55), Inches(1.6), Inches(3.9), "≈ 86 %", "Progression du sprint", BLEU)
chiffre(s, Inches(8.65), Inches(1.6), Inches(4.2), "3", "Tickets transférés au Sprint 3", GRIS)
carte(s, Inches(0.45), Inches(3.35), Inches(8.5), Inches(3.4),
      "Répartition des responsabilités (tâches terminées)",
      ["Souleymane DIALLO — authentification (15,16,17), réservation + anti-double (21),",
       "        squelettes du monorepo (28), Docker Compose (29), intégration",
       "Zakaria Lahouiri — disponibilités des médecins (20), flux clinique du jour (23)",
       "Larbi Saib — patient léger (35), réservation par la réception (36), socle RDV (49)",
       "Frontend d'authentification & refonte UI (42–48) — terminés, effort partagé",
       " ",
       "En cours (Sprint 3) : notifications internes (25, Zakaria)",
       "Livré en avance sur le Sprint 3 : annulation d'un RDV avec motif (22)"], BLEU)
carte(s, Inches(9.15), Inches(3.35), Inches(3.7), Inches(3.4),
      "Transféré au Sprint 3",
      ["CRUD médecins par l'admin",
       "   (MEDIPLAN-34)",
       "Précision doc UC-02",
       "   (MEDIPLAN-38)",
       "Secrets / variables",
       "   d'environnement (MEDIPLAN-39)",
       " ",
       "Non bloquants pour la",
       "   démonstration"], GRIS)
pied(s, "Souleymane", "1 min")
notes(s, "SOULEYMANE — 1 min — NOMMER CHAQUE PERSONNE À VOIX HAUTE (exigence de la "
         "consigne).\n\n"
         "« 18 tickets terminés sur 21, environ 86 % du sprint. »\n\n"
         "« La répartition : j'ai porté l'authentification, la réservation avec "
         "l'anti-double-réservation, la mise en place du monorepo et de Docker, et "
         "l'intégration. Zakaria a fait les disponibilités des médecins et le flux clinique "
         "du jour. Larbi a modélisé le patient léger, la réservation par la réception et le "
         "socle technique des rendez-vous. Le frontend d'authentification et la refonte de "
         "l'interface, c'est un effort partagé. »\n\n"
         "« Trois tickets passent au Sprint 3 : la gestion des médecins par "
         "l'administrateur, une précision de documentation, et les variables "
         "d'environnement. Aucun n'est bloquant, et nous les assumons comme tels. Et nous "
         "avons même livré l'annulation en avance sur le planning. »\n\n"
         "→ TRANSITION : « Zakaria, pour la suite. »")

# =========================================================================
# DIAPO 10 — Prochaines etapes / Sprint 3
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 9, "Prochaines étapes — Sprint 3", "Ce qui reste et nos priorités")
carte(s, Inches(0.45), Inches(1.65), Inches(6.1), Inches(3.1),
      "Priorités du Sprint 3 (cycle de vie du RDV)",
      ["Annuler un RDV avec motif — déjà livré, à valider (22)",
       "Notifications internes sur les changements (25) — en cours",
       "Statistiques complètes du tableau de bord (26)",
       "Sécurité RDV : garde-fou patient léger + accès",
       "   en lecture cloisonné (50)",
       "Décaler en bloc les rendez-vous d'un médecin (24)"], SARCELLE)
carte(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(3.1),
      "Reporté depuis le Sprint 2",
      ["Gestion des médecins par l'admin (34)",
       "Précision documentaire UC-02 (38)",
       "Variables d'environnement & secrets (39)",
       " ",
       "Puis, en finalisation :",
       "   export CSV, intégration continue (CI),",
       "   déploiement, documentation des tests"], GRIS)
carte(s, Inches(0.45), Inches(4.95), Inches(12.4), Inches(1.8),
      "En un mot",
      ["Le Sprint 2 livre une tranche complète : se connecter → définir des disponibilités → réserver → suivre la journée.",
       "Le Sprint 3 ajoute le cycle de vie (annulation, notifications, statistiques) et durcit la sécurité des rendez-vous."],
      BLEU)
pied(s, "Zakaria", "1 min")
notes(s, "ZAKARIA — 1 min\n\n"
         "« Le Sprint 3 porte sur le cycle de vie complet du rendez-vous : l'annulation, "
         "déjà livrée, à valider officiellement ; les notifications internes, en cours ; les "
         "statistiques réelles du tableau de bord ; et un renforcement de la sécurité des "
         "rendez-vous. »\n\n"
         "« En un mot : le Sprint 2 livre une tranche complète, de la connexion au suivi de "
         "la journée. Le Sprint 3 ajoute le cycle de vie et durcit la sécurité. »\n\n"
         "PRÉCISION SI ON DEMANDE : le compteur des rendez-vous du jour est déjà réel ; ce "
         "sont les statistiques complètes qui viennent au Sprint 3.\n\n"
         "NE PAS PROMETTRE le décalage en bloc (24) : le code existe sur une branche, il "
         "n'est pas intégré.")

# =========================================================================
# DIAPO 11 — Merci / Questions
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(3.5), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.4))
para(tf, "Merci de votre attention", taille=40, couleur=BLANC, gras=True,
     premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(3.8), Inches(11.3), Inches(1.6))
para(tf, "Questions ?", taille=28, couleur=BLEU_PALE, gras=True,
     premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
para(tf, "Équipe MediPlan — Souleymane DIALLO  ·  Zakaria Lahouiri  ·  Larbi Saib",
     taille=16, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=14)
notes(s, "SOULEYMANE — 15 s — REPRENDRE L'ACCROCHE DU DÉBUT. C'est ce qui donne une "
         "présentation construite plutôt que trois exposés.\n\n"
         "« On a commencé par un créneau vendu deux fois. Aujourd'hui, dans MediPlan, c'est "
         "la base de données elle-même qui l'empêche. La réception se connecte, ouvre son "
         "agenda, réserve, suit sa journée, annule — de bout en bout, avec 171 tests "
         "derrière. Merci. Nous répondons à vos questions, chacun sur sa partie. »\n\n"
         "RÈGLE : répondre court, puis SE TAIRE. Une réponse qui s'étire donne l'impression "
         "qu'on cherche.\n\n"
         "— Pourquoi le patient ne réserve pas lui-même ? (Souleymane) « Choix métier "
         "assumé : l'utilisateur réel est la réception, le patient appelle. Le libre-service "
         "viendra quand le cœur sera solide. »\n"
         "— Comment évitez-vous les doubles réservations ? (Souleymane) « Un index unique "
         "partiel en base, plus une transaction avec verrou. Ce n'est pas notre code qui "
         "arbitre, c'est PostgreSQL. Testé en concurrence. »\n"
         "— Les données de santé sont-elles protégées ? (Souleymane) « Quatre rôles, "
         "cloisonnement par clinique, mots de passe hachés bcrypt, verrouillage après cinq "
         "échecs, et un patient léger qui ne peut pas se connecter par construction. »\n"
         "— Qu'est-ce qui n'est pas fini ? (Souleymane) « Trois tickets non bloquants "
         "transférés au Sprint 3 : gestion des médecins, une précision de doc, les secrets "
         "d'environnement. »\n"
         "— Comment testez-vous le frontend ? (Larbi) « Composants, services HTTP, guards de "
         "rôle et intercepteurs — 123 tests — plus une vérification manuelle bout-en-bout. »\n"
         "— Pourquoi Angular / NestJS ? (Zakaria) « Un seul langage, TypeScript, du "
         "navigateur à la base. Une équipe de trois, une seule courbe d'apprentissage. »\n"
         "— Et le décalage en bloc ? (Souleymane) « Le code existe sur une branche, il n'est "
         "pas intégré. Donc pour nous il n'est pas terminé — c'est le Sprint 3. »\n"
         "— Qui a fait quoi ? Chacun reprend SA ligne de la diapo Jira.\n\n"
         "SI ON NE SAIT PAS : « Je ne veux pas vous répondre de mémoire — c'est un point que "
         "je vérifie et je vous reviens dessus. » Jamais de bluff.")

# --- Sauvegarde ----------------------------------------------------------
sortie = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "MediPlan-Sprint2-Presentation.pptx")
prs.save(sortie)
print("Genere :", sortie)
print("Nombre de diapos :", len(prs.slides._sldIdLst))
