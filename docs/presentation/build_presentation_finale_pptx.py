# -*- coding: utf-8 -*-
"""Genere le support de la presentation finale de MediPlan (PowerPoint).

Presentation finale du 13 aout 2026 — 15 minutes + 5 minutes de questions.
Structure calquee sur les 7 sections de la consigne, repartie sur 3 orateurs.

Usage : python docs/presentation/build_presentation_finale_pptx.py
Sortie : docs/presentation/MediPlan-Presentation-Finale.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette (identique aux supports precedents, pour la continuite) ---
BLEU = RGBColor(0x15, 0x4D, 0x8C)
BLEU_CLAIR = RGBColor(0x1E, 0x88, 0xE5)
SARCELLE = RGBColor(0x00, 0x89, 0x7B)
GRIS = RGBColor(0x40, 0x40, 0x40)
GRIS_CLAIR = RGBColor(0x6B, 0x6B, 0x6B)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
FOND = RGBColor(0xF4, 0xF7, 0xFB)
BLEU_PALE = RGBColor(0xCF, 0xE3, 0xFF)
BLEU_PALE2 = RGBColor(0x9F, 0xC8, 0xF5)
AMBRE = RGBColor(0xB5, 0x6A, 0x00)
ROUGE = RGBColor(0xA8, 0x32, 0x32)

LARGEUR = Inches(13.333)
HAUTEUR = Inches(7.5)

prs = Presentation()
prs.slide_width = LARGEUR
prs.slide_height = HAUTEUR
BLANK = prs.slide_layouts[6]


def fond(slide, couleur=BLANC):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = couleur


def rect(slide, x, y, w, h, couleur):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = couleur
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def zone_texte(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def para(tf, texte, taille=18, couleur=GRIS, gras=False, puce=False,
         espace_avant=6, niveau=0, aligne=PP_ALIGN.LEFT, premier=False):
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    p.alignment = aligne
    p.level = niveau
    p.space_after = Pt(espace_avant)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = (("•  " if puce else "") + texte)
    run.font.size = Pt(taille)
    run.font.bold = gras
    run.font.color.rgb = couleur
    run.font.name = "Calibri"
    return p


def notes(slide, texte):
    slide.notes_slide.notes_text_frame.text = texte


def bandeau_titre(slide, numero, titre, sous_titre=None):
    rect(slide, 0, 0, LARGEUR, Inches(1.25), BLEU)
    rect(slide, 0, Inches(1.25), LARGEUR, Pt(6), SARCELLE)
    past = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.32),
                                  Inches(0.62), Inches(0.62))
    past.fill.solid(); past.fill.fore_color.rgb = SARCELLE
    past.line.fill.background(); past.shadow.inherit = False
    pf = past.text_frame; pf.word_wrap = False
    pp = pf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); r.text = str(numero); r.font.size = Pt(22)
    r.font.bold = True; r.font.color.rgb = BLANC
    tb, tf = zone_texte(slide, Inches(1.25), Inches(0.18), Inches(11.6), Inches(1.0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, titre, taille=28, couleur=BLANC, gras=True, premier=True, espace_avant=0)
    if sous_titre:
        para(tf, sous_titre, taille=14, couleur=BLEU_PALE, espace_avant=0)


def pied(slide, orateur, duree):
    tb, tf = zone_texte(slide, Inches(0.45), Inches(7.0), Inches(12.4), Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = f"Orateur : {orateur}"
    r.font.size = Pt(11); r.font.color.rgb = GRIS_CLAIR; r.font.italic = True
    tb2, tf2 = zone_texte(slide, Inches(10.0), Inches(7.0), Inches(2.9), Inches(0.4))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"⏱  {duree}"
    r2.font.size = Pt(11); r2.font.color.rgb = SARCELLE; r2.font.bold = True


def carte(slide, x, y, w, h, titre, lignes, couleur_titre=BLEU):
    rect(slide, x, y, w, h, FOND)
    rect(slide, x, y, Pt(6), h, couleur_titre)
    tb, tf = zone_texte(slide, x + Inches(0.25), y + Inches(0.12),
                        w - Inches(0.4), h - Inches(0.2))
    para(tf, titre, taille=16, couleur=couleur_titre, gras=True, premier=True,
         espace_avant=4)
    for l in lignes:
        para(tf, l, taille=13, couleur=GRIS, puce=bool(l.strip()), espace_avant=3)


def chiffre(slide, x, y, largeur, valeur, libelle, couleur):
    rect(slide, x, y, largeur, Inches(1.5), couleur)
    tb, tf = zone_texte(slide, x, y + Inches(0.08), largeur, Inches(1.4))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, valeur, taille=40, couleur=BLANC, gras=True, premier=True,
         aligne=PP_ALIGN.CENTER, espace_avant=0)
    para(tf, libelle, taille=13, couleur=BLANC, aligne=PP_ALIGN.CENTER, espace_avant=0)


# =========================================================================
# DIAPO 1 — Titre
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(3.05), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(1.3), Inches(11.3), Inches(1.4))
para(tf, "MediPlan", taille=60, couleur=BLANC, gras=True, premier=True,
     aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(2.25), Inches(11.3), Inches(0.7))
para(tf, "Plateforme web de gestion des rendez-vous médicaux", taille=22,
     couleur=BLEU_PALE, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(3.3), Inches(11.3), Inches(1.4))
para(tf, "Présentation finale", taille=30, couleur=BLANC,
     gras=True, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=8)
para(tf, "Une solution déployée, testée et en ligne", taille=18,
     couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.5))
para(tf, "Souleymane DIALLO  ·  Zakaria Lahouiri  ·  Larbi Saib", taille=18,
     couleur=BLANC, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=4)
para(tf, "Projet intégrateur 030747  ·  Collège La Cité  ·  13 août 2026",
     taille=14, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=2)
notes(s, "SOULEYMANE — 30 secondes.\n\n"
         "« Bonjour. Nous sommes Souleymane, Zakaria et Larbi, et nous vous présentons "
         "MediPlan : une plateforme web de gestion de rendez-vous pour une clinique "
         "médicale. La solution est en ligne, nous vous la montrerons en direct.\n\n"
         "Nous avons quinze minutes. Je commence par le problème et la solution, Larbi "
         "enchaîne sur notre démarche, je reviens sur l'architecture, Zakaria fait la "
         "démonstration, puis nous terminons sur les tests, les limites et ce que chacun "
         "a fait. »\n\n"
         "NE PAS lire la diapo. Regarder la salle.")

# =========================================================================
# DIAPO 2 — Problematique
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 1, "Le problème",
              "Ce que vit une petite clinique qui gère ses rendez-vous au téléphone")
carte(s, Inches(0.45), Inches(1.65), Inches(6.1), Inches(4.9),
      "Le quotidien de la réception",
      ["Le téléphone est le canal principal — le patient appelle,",
       "   il ne s'inscrit pas en ligne",
       "",
       "Un agenda papier ou un tableur partagé : deux personnes",
       "   peuvent réserver le même créneau au même moment",
       "",
       "Aucune vue d'ensemble de la journée en cours : qui est",
       "   arrivé, qui est en consultation, qui n'est pas venu",
       "",
       "Les absences (no-show) ne sont ni tracées ni mesurées",
       "",
       "Le médecin et la réception ne partagent pas la même",
       "   information au même instant"], ROUGE)
carte(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(4.9),
      "Ce que ça coûte",
      ["Des doubles réservations à gérer en catastrophe,",
       "   devant le patient",
       "",
       "Du temps perdu à chercher l'information plutôt qu'à",
       "   accueillir",
       "",
       "Des créneaux perdus : un rendez-vous annulé ne redevient",
       "   jamais disponible",
       "",
       "Aucune donnée pour décider — impossible de savoir si",
       "   l'agenda est bien rempli ou si les absences explosent",
       "",
       "Une charge mentale portée par une seule personne"], AMBRE)
pied(s, "Souleymane DIALLO", "1 min")
notes(s, "SOULEYMANE — 1 minute.\n\n"
         "« Le point de départ, c'est une petite clinique médicale. La réception gère les "
         "rendez-vous au téléphone, avec un agenda papier ou un tableur.\n\n"
         "Trois irritants reviennent. D'abord la double réservation : deux personnes "
         "peuvent inscrire un patient sur le même créneau. Ensuite, personne n'a de vue "
         "sur la journée en cours — qui est arrivé, qui est en consultation. Et enfin, "
         "quand un rendez-vous est annulé, le créneau est perdu : il ne revient jamais "
         "dans les disponibilités.\n\n"
         "Ce ne sont pas des problèmes théoriques. Ce sont ceux que nous avons retenus "
         "pour cadrer le produit. »\n\n"
         "Insister sur la double réservation : elle revient en démonstration.")

# =========================================================================
# DIAPO 3 — Solution, utilisateurs, objectifs
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 2, "Notre solution",
              "Outiller la réception, pas remplacer le téléphone")
carte(s, Inches(0.45), Inches(1.65), Inches(4.0), Inches(4.9),
      "Ce que fait MediPlan",
      ["Le médecin publie ses plages de",
       "   disponibilité ; le système génère",
       "   les créneaux réservables tout seul",
       "",
       "La réception réserve pour un patient",
       "   au téléphone, en quelques secondes",
       "",
       "La base garantit qu'un créneau ne",
       "   peut être pris deux fois",
       "",
       "Le flux du jour suit chaque patient :",
       "   arrivé, en consultation, terminé",
       "",
       "Annulation avec motif — le créneau",
       "   redevient disponible"], SARCELLE)
carte(s, Inches(4.65), Inches(1.65), Inches(4.0), Inches(4.9),
      "Pour qui",
      ["Réception / administrateur de",
       "   clinique — l'utilisateur principal,",
       "   celui qui vit le problème",
       "",
       "Médecin — consulte son agenda et",
       "   le flux de sa journée",
       "",
       "Patient — existe dans le système",
       "   comme « patient léger », créé au",
       "   comptoir, sans compte à retenir",
       "",
       "Super administrateur — supervise",
       "   l'ensemble des cliniques"], BLEU)
carte(s, Inches(8.85), Inches(1.65), Inches(4.0), Inches(4.9),
      "Nos objectifs",
      ["Rendre la double réservation",
       "   techniquement impossible",
       "",
       "Donner une vue partagée et fiable",
       "   de la journée en cours",
       "",
       "Ne perdre aucun créneau annulé",
       "",
       "Mesurer l'activité : occupation,",
       "   absences, volumes par médecin",
       "",
       "Livrer une application réellement",
       "   déployée, pas une maquette"], BLEU_CLAIR)
pied(s, "Souleymane DIALLO", "1 min")
notes(s, "SOULEYMANE — 1 minute.\n\n"
         "« Notre parti pris : nous n'essayons pas de remplacer le téléphone. Dans une "
         "petite clinique, le patient appelle, c'est le canal réel. Nous outillons la "
         "personne qui répond.\n\n"
         "D'où une décision importante : le patient léger. La réception crée le patient au "
         "comptoir, avec un nom et un téléphone, sans compte ni mot de passe. Le patient "
         "n'a rien à retenir. C'est le flux réel d'une clinique.\n\n"
         "Nos objectifs, ce sont les irritants d'avant, retournés en critères vérifiables. "
         "Le premier — rendre la double réservation impossible — je vous montrerai comment "
         "nous l'avons tenu techniquement. »\n\n"
         "Si on demande : « pourquoi le patient ne réserve pas lui-même ? » → « le libre-"
         "service existe pour l'inscription patient, mais la prise de RDV passe par la "
         "réception : c'est ce que fait réellement une petite clinique ».")

# =========================================================================
# DIAPO 4 — Organisation de l'equipe
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 3, "Comment nous avons travaillé",
              "Trois personnes, des responsabilités explicites")
carte(s, Inches(0.45), Inches(1.65), Inches(4.0), Inches(4.4),
      "Souleymane DIALLO",
      ["Pilotage, Jira et GitHub",
       "Cahier des charges,",
       "   7 cas d'utilisation",
       "Authentification et RBAC",
       "Socle monorepo + Docker",
       "Design system et refonte UI",
       "Déploiement Azure (Bicep)"], BLEU)
carte(s, Inches(4.65), Inches(1.65), Inches(4.0), Inches(4.4),
      "Zakaria Lahouiri",
      ["Diagramme de classes",
       "Disponibilités des médecins",
       "   et génération des créneaux",
       "Flux clinique du jour",
       "Notifications internes",
       "Export CSV des rendez-vous"], SARCELLE)
carte(s, Inches(8.85), Inches(1.65), Inches(4.0), Inches(4.4),
      "Larbi Saib",
      ["Diagrammes de séquence",
       "Modèle du « patient léger »",
       "Prise de RDV par la réception",
       "Socle technique des RDV et",
       "   index anti-double-réservation",
       "Tableau de bord et statistiques"], BLEU_CLAIR)
rect(s, Inches(0.45), Inches(6.2), Inches(12.4), Inches(0.62), FOND)
tb, tf = zone_texte(s, Inches(0.7), Inches(6.27), Inches(12.0), Inches(0.5))
para(tf, "Découpage par tranche verticale : chacun mène sa fonctionnalité de la base de "
         "données jusqu'à l'écran — pas de séparation « un front / un back ».",
     taille=13, couleur=BLEU, gras=True, premier=True, espace_avant=0)
pied(s, "Larbi Saib", "1 min")
notes(s, "LARBI — 1 minute.\n\n"
         "« Nous sommes trois, et nous avons réparti par fonctionnalité plutôt que par "
         "couche technique. Chacun mène sa fonctionnalité de bout en bout : la migration "
         "de base de données, l'API, l'écran Angular et les tests.\n\n"
         "Ce choix a un avantage et un coût. L'avantage : personne n'attend après personne "
         "pour livrer quelque chose de visible. Le coût : il faut se coordonner sur les "
         "fichiers partagés, et nous y reviendrons dans les difficultés.\n\n"
         "Concrètement, j'ai porté le patient léger, la prise de rendez-vous par la "
         "réception et le tableau de bord. Zakaria a porté les disponibilités, le flux du "
         "jour, les notifications et l'export. Souleymane a porté l'authentification, le "
         "socle technique et le déploiement. »")

# =========================================================================
# DIAPO 5 — Jira et GitHub
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 4, "Jira et GitHub", "Notre flux de travail, et ce qu'il nous a appris")
chiffre(s, Inches(0.45), Inches(1.6), Inches(2.85), "51", "Tickets Jira", BLEU)
chiffre(s, Inches(3.55), Inches(1.6), Inches(2.85), "21", "Pull requests", SARCELLE)
chiffre(s, Inches(6.65), Inches(1.6), Inches(2.85), "3", "Contributeurs sur main", BLEU_CLAIR)
chiffre(s, Inches(9.75), Inches(1.6), Inches(2.85), "7", "Épiques", GRIS)
carte(s, Inches(0.45), Inches(3.35), Inches(6.1), Inches(3.3),
      "Notre flux",
      ["Un ticket Jira → une branche → une pull request → une revue",
       "La CI GitHub Actions bloque la fusion si le build ou les",
       "   tests échouent",
       "Chaque ticket livré porte en commentaire le lien de sa PR",
       "« Terminé » veut dire fusionné dans main — pas « codé »"], BLEU)
carte(s, Inches(6.75), Inches(3.35), Inches(6.1), Inches(3.3),
      "Notre plus grande difficulté : la dérive d'intégration",
      ["Des branches sont restées non fusionnées pendant des",
       "   semaines, jusqu'à 56 commits de retard",
       "Cinq tickets étaient marqués « Terminé » sans exister",
       "   dans le produit",
       "Correction : trois branches réintégrées en trois PR revues,",
       "   les autres repassées « À faire » avec justification",
       "La leçon : un ticket n'est terminé qu'une fois fusionné"], AMBRE)
pied(s, "Larbi Saib", "1 min")
notes(s, "LARBI — 1 minute. C'EST LA DIAPO LA PLUS IMPORTANTE DE MA PARTIE.\n\n"
         "« Notre flux est classique : un ticket, une branche, une pull request, une revue. "
         "La CI bloque si les tests cassent.\n\n"
         "Mais je veux être honnête sur notre principale difficulté, parce que c'est ce "
         "qui nous a le plus appris. Nous avons travaillé chacun sur nos branches sans les "
         "fusionner assez souvent. Résultat : des branches à cinquante commits de retard, "
         "et des tickets marqués Terminé dans Jira alors que le code n'était nulle part "
         "dans le produit.\n\n"
         "Nous l'avons corrigé : trois fonctionnalités ont été réintégrées en trois pull "
         "requests revues, et les tickets dont le code n'était pas fusionné sont repassés "
         "À faire, avec la raison écrite dans Jira.\n\n"
         "Ce que nous en retenons : dans notre définition de terminé, terminé veut dire "
         "fusionné dans main. Pas codé sur une branche. »\n\n"
         "ASSUMER. Ne pas s'excuser. C'est une leçon de génie logiciel, pas une faute.")

# =========================================================================
# DIAPO 6 — Architecture
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 5, "Architecture générale",
              "Un monorepo, deux applications, une base — et le tout en ligne")
carte(s, Inches(0.45), Inches(1.65), Inches(3.9), Inches(3.1),
      "Frontend — Angular 22",
      ["Composants standalone, Signals",
       "Angular Material 3 + Tailwind 4",
       "Design system : une seule source",
       "   de tokens, mode sombre inclus",
       "Garde de rôle sur chaque route"], BLEU_CLAIR)
carte(s, Inches(4.55), Inches(1.65), Inches(3.9), Inches(3.1),
      "Backend — NestJS 11",
      ["API REST préfixée api/v1",
       "TypeORM, validation par DTO",
       "JWT + gardes RBAC à 4 rôles",
       "Chaque requête bornée à la",
       "   clinique de l'appelant"], SARCELLE)
carte(s, Inches(8.65), Inches(1.65), Inches(4.2), Inches(3.1),
      "Données — PostgreSQL",
      ["Schéma piloté uniquement par",
       "   migrations versionnées",
       "Aucune synchronisation",
       "   automatique du schéma",
       "Index unique partiel qui interdit",
       "   la double réservation"], BLEU)
carte(s, Inches(0.45), Inches(5.0), Inches(12.4), Inches(1.65),
      "Comment les parties se parlent — et pourquoi c'est sûr",
      ["Le navigateur ne parle qu'au frontend. Le frontend relaie les appels /api vers le backend.",
       "Le backend n'a AUCUNE adresse publique (ingress interne) : il n'est joignable que par le frontend — ce qui supprime tout besoin de CORS.",
       "Hébergement : Azure Container Apps, base PostgreSQL infogérée chez Neon, images sur ghcr.io. Toute l'infrastructure est décrite en Bicep. Coût : environ 0 $/mois."], BLEU)
pied(s, "Souleymane DIALLO", "1 min 30")
notes(s, "SOULEYMANE — 1 minute 30.\n\n"
         "« L'architecture tient en une phrase : un monorepo avec deux applications et une "
         "base de données.\n\n"
         "Trois choix méritent d'être expliqués.\n\n"
         "Premier choix : le schéma de la base n'est piloté que par des migrations "
         "versionnées. Nous avons désactivé la synchronisation automatique de TypeORM. "
         "C'est plus contraignant, mais le schéma est reproductible : n'importe qui "
         "reconstruit exactement la même base.\n\n"
         "Deuxième choix : le backend n'a pas d'adresse publique. Il est en ingress "
         "interne, seul le frontend peut l'appeler. Conséquence directe : nous n'avons "
         "aucune configuration CORS à gérer, parce qu'il n'y a jamais de requête "
         "inter-origines.\n\n"
         "Troisième choix : toute l'infrastructure est décrite en code, en Bicep. Aucune "
         "ressource n'a été créée à la main dans le portail Azure. Et grâce au "
         "scale-to-zero, ça nous coûte environ zéro dollar par mois — ce qui comptait, "
         "notre crédit étudiant n'était pas renouvelable. »")

# =========================================================================
# DIAPO 7 — La decision technique dont on est le plus fier
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 6, "La double réservation",
              "Le choix technique dont nous sommes le plus fiers")
carte(s, Inches(0.45), Inches(1.65), Inches(6.1), Inches(2.5),
      "L'approche naïve — et pourquoi elle échoue",
      ["« Je vérifie si le créneau est libre, puis je réserve »",
       "Entre la vérification et l'écriture, une autre requête peut",
       "   passer : les deux voient le créneau libre",
       "Deux réceptionnistes qui cliquent en même temps créent",
       "   deux rendez-vous sur le même créneau"], ROUGE)
carte(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(2.5),
      "Notre solution — c'est la base qui tranche",
      ["Un index unique partiel en PostgreSQL sur le créneau,",
       "   qui ignore les rendez-vous annulés",
       "La base refuse physiquement la seconde écriture",
       "Ce n'est plus le code qui arbitre la course, c'est",
       "   PostgreSQL — et lui ne se trompe jamais"], SARCELLE)
carte(s, Inches(0.45), Inches(4.4), Inches(12.4), Inches(2.2),
      "Ce que ce choix nous donne en plus",
      ["L'annulation est « hors index » : un rendez-vous annulé ne bloque plus le créneau, qui redevient donc réservable. Le deuxième irritant est réglé par le même mécanisme.",
       "Vérifié sur l'application déployée : deux réservations lancées en même temps sur le même créneau — une passe (201), l'autre est refusée (409).",
       "La garantie ne dépend pas de la qualité du code applicatif : même si une nouvelle fonctionnalité oubliait la vérification, la base tiendrait."], BLEU)
pied(s, "Souleymane DIALLO", "1 min")
notes(s, "SOULEYMANE — 1 minute. Diapo à ne PAS survoler : c'est notre meilleur argument "
         "technique, et la question tombe souvent.\n\n"
         "« Notre premier objectif était de rendre la double réservation impossible. "
         "L'approche naïve serait : je vérifie que le créneau est libre, puis j'écris. Le "
         "problème, c'est qu'entre les deux, une autre requête peut passer. Les deux voient "
         "le créneau libre, les deux réservent.\n\n"
         "Notre solution, c'est de ne pas laisser le code arbitrer. Nous avons posé un "
         "index unique partiel en PostgreSQL : partiel, parce qu'il ignore les rendez-vous "
         "annulés. La base refuse physiquement la deuxième écriture.\n\n"
         "Et ce même index nous donne l'annulation gratuitement : comme un rendez-vous "
         "annulé sort de l'index, il ne bloque plus le créneau, qui redevient réservable. "
         "Un seul mécanisme règle deux irritants. »\n\n"
         "Si on demande « et si le code oublie la vérification ? » → « c'est exactement "
         "l'intérêt : la garantie est dans la base, pas dans la discipline du développeur ».")

# =========================================================================
# DIAPO 8 — Demonstration (annonce du scenario)
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 7, "Démonstration", "Une journée à la réception, du matin au soir")
carte(s, Inches(0.45), Inches(1.65), Inches(6.1), Inches(4.6),
      "Le scénario",
      ["1.  Je me connecte comme la réception de la clinique",
       "2.  Le tableau de bord m'affiche l'activité réelle",
       "3.  Le Dr Bergeron publie une plage du matin —",
       "     les créneaux se génèrent automatiquement",
       "4.  Un patient appelle : je le crée et je réserve",
       "5.  Le flux du jour : arrivé → en consultation → terminé",
       "6.  Un patient annule : motif obligatoire,",
       "     le créneau redevient libre",
       "7.  J'exporte les rendez-vous du mois en CSV",
       "8.  Je bascule sur le compte médecin :",
       "     le menu change, et sa cloche a sonné"], SARCELLE)
carte(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(4.6),
      "Ce que vous verrez au passage",
      ["Une application réellement en ligne, sur Azure —",
       "   pas un serveur local",
       "",
       "Le contrôle d'accès visible à l'écran : le médecin",
       "   n'a tout simplement pas les mêmes entrées de menu",
       "",
       "Les notifications qui se déclenchent toutes seules",
       "   quand un rendez-vous change d'état",
       "",
       "Le mode sombre complet, sans qu'aucun écran",
       "   n'ait été retouché à la main",
       "",
       "Des statistiques calculées sur de vraies données"], BLEU)
pied(s, "Zakaria Lahouiri", "30 s + 4 min de démo")
notes(s, "ZAKARIA — 30 secondes, puis on bascule sur le navigateur.\n\n"
         "« Je vous propose de suivre une journée à la réception, du matin au soir. Huit "
         "étapes, et vous verrez au passage le contrôle d'accès, les notifications et les "
         "statistiques.\n\n"
         "Un mot important : ce que je vais vous montrer tourne en ligne, sur Azure. Ce "
         "n'est pas un serveur sur mon portable. »\n\n"
         "AVANT DE PASSER : l'application doit avoir été réveillée 15 minutes plus tôt "
         "(scale-to-zero). Onglet déjà ouvert sur l'écran de connexion, session fermée.\n"
         "SI LA DÉMO ÉCHOUE : passer immédiatement à la vidéo, sans commenter l'incident.")

# =========================================================================
# DIAPO 9 — Ecran de secours pendant la demo
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(3.15), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.2))
para(tf, "Démonstration en direct", taille=44, couleur=BLANC, gras=True,
     premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(2.4))
para(tf, "ca-mediplan-frontend.ashytree-9ad5012f.canadacentral.azurecontainerapps.io",
     taille=16, couleur=BLEU_PALE, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=14)
para(tf, "Réception : admin.demo@mediplan.test    ·    Médecin : doctor.demo@mediplan.test",
     taille=15, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=8)
para(tf, "Comptes et données de démonstration uniquement",
     taille=12, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=6)
notes(s, "ZAKARIA — 4 minutes de démonstration. Cette diapo reste affichée si on revient "
         "au support.\n\n"
         "DÉROULÉ EXACT :\n"
         "1. Connexion admin.demo / Adm1n!Secret — « je suis la réception »\n"
         "2. Tableau de bord — pointer les chiffres réels\n"
         "3. Disponibilités → créer une plage → ouvrir les créneaux générés\n"
         "4. Nouveau rendez-vous → créer le patient → réserver\n"
         "5. Flux du jour → Arrive → Consultation → Termine\n"
         "6. Menu ⋯ → Annuler → motif obligatoire\n"
         "7. Export CSV → ouvrir le fichier téléchargé\n"
         "8. Déconnexion → doctor.demo / Doct0r!Secret → montrer le menu réduit ET la cloche\n\n"
         "NE JAMAIS FAIRE F5 pendant la démonstration.\n"
         "Parler pendant les chargements, ne pas laisser de silence.")

# =========================================================================
# DIAPO 10 — Valeur ajoutee
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 8, "Ce que ça change pour la clinique",
              "Chaque irritant du départ, et la réponse apportée")
y = Inches(1.7)
lignes = [
    ("Deux personnes réservaient le même créneau",
     "La base refuse la seconde écriture — ce n'est plus une question de vigilance", SARCELLE),
    ("Un créneau annulé était perdu",
     "L'annulation libère le créneau, qui redevient immédiatement réservable", SARCELLE),
    ("Personne n'avait de vue sur la journée en cours",
     "Le flux du jour est partagé : réception et médecin voient le même état", BLEU_CLAIR),
    ("Le patient devait se créer un compte pour exister",
     "Le patient léger est créé au comptoir, sans mot de passe à retenir", BLEU_CLAIR),
    ("Aucune donnée pour piloter l'activité",
     "Taux d'occupation, taux d'absence et volumes par médecin, calculés en direct", BLEU),
    ("Le médecin apprenait les changements de vive voix",
     "Une notification interne est émise à chaque réservation, changement ou annulation", BLEU),
]
for avant, apres, coul in lignes:
    rect(s, Inches(0.45), y, Inches(5.9), Inches(0.72), FOND)
    rect(s, Inches(0.45), y, Pt(5), Inches(0.72), ROUGE)
    tb, tf = zone_texte(s, Inches(0.68), y + Inches(0.08), Inches(5.6), Inches(0.6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, avant, taille=12, couleur=GRIS_CLAIR, premier=True, espace_avant=0)
    fl = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.48), y + Inches(0.22),
                            Inches(0.42), Inches(0.28))
    fl.fill.solid(); fl.fill.fore_color.rgb = coul
    fl.line.fill.background(); fl.shadow.inherit = False
    rect(s, Inches(7.05), y, Inches(5.8), Inches(0.72), FOND)
    rect(s, Inches(7.05), y, Pt(5), Inches(0.72), coul)
    tb, tf = zone_texte(s, Inches(7.28), y + Inches(0.08), Inches(5.5), Inches(0.6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, apres, taille=12, couleur=GRIS, gras=True, premier=True, espace_avant=0)
    y += Inches(0.8)
pied(s, "Zakaria Lahouiri", "45 s")
notes(s, "ZAKARIA — 45 secondes. Ne PAS lire les six lignes.\n\n"
         "« Je reviens une seconde sur le problème du début. À gauche, les irritants que "
         "Souleymane a présentés. À droite, ce que la solution en fait.\n\n"
         "Je n'en commente que deux. La double réservation n'est plus une question de "
         "vigilance humaine : c'est la base qui refuse. Et le créneau annulé n'est plus "
         "perdu : il redevient disponible tout de suite, ce qui est du chiffre d'affaires "
         "récupéré pour la clinique.\n\n"
         "Le reste suit la même logique : chaque irritant du départ a une réponse "
         "dans le produit. »")

# =========================================================================
# DIAPO 11 — Tests et validation
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 9, "Tests et validation", "Ce que nous vérifions, et à quel moment")
chiffre(s, Inches(0.45), Inches(1.6), Inches(2.85), "186", "Tests automatisés verts", SARCELLE)
chiffre(s, Inches(3.55), Inches(1.6), Inches(2.85), "57", "Tests backend", BLEU)
chiffre(s, Inches(6.65), Inches(1.6), Inches(2.85), "129", "Tests frontend", BLEU_CLAIR)
chiffre(s, Inches(9.75), Inches(1.6), Inches(2.85), "0", "Fusion sans CI verte", GRIS)
carte(s, Inches(0.45), Inches(3.35), Inches(6.1), Inches(3.3),
      "Notre stratégie de test",
      ["Backend : les règles métier et la sécurité — transitions de",
       "   statut, droits par rôle, périmètre de clinique",
       "Le cas critique vérifié en concurrence sur l'application en",
       "   ligne : deux réservations simultanées, une seule passe.",
       "   Vérification manuelle et reproductible — nos tests backend",
       "   simulent la base, ils ne peuvent pas le prouver seuls",
       "Frontend : le comportement des écrans et des formulaires,",
       "   pas leur apparence",
       "Validation manuelle : le scénario complet rejoué sur",
       "   l'application en ligne avant chaque démonstration"], BLEU)
carte(s, Inches(6.75), Inches(3.35), Inches(6.1), Inches(3.3),
      "Ce que la CI garantit à chaque push",
      ["Le projet compile — backend et frontend",
       "Les 186 tests passent",
       "Les templates d'infrastructure Bicep sont valides",
       "Une pull request dont la CI est rouge n'est pas fusionnée",
       "Le lint et le formatage sont rapportés sans bloquer :",
       "   un choix assumé, pour ne pas avoir une CI rouge en",
       "   permanence, donc ignorée"], SARCELLE)
pied(s, "Larbi Saib", "1 min")
notes(s, "LARBI — 1 minute.\n\n"
         "« Nous avons 186 tests automatisés : 57 côté serveur, 129 côté interface. Ils "
         "tournent à chaque push, et une pull request dont la CI est rouge n'est pas "
         "fusionnée.\n\n"
         "La vérification dont je suis le plus content n'est pas dans cette suite, et je "
         "veux être précis là-dessus. Pour la double réservation, nous avons lancé deux "
         "réservations en même temps sur le même créneau, sur l'application en ligne : une "
         "passe, l'autre est refusée avec un 409. C'est reproductible, mais c'est manuel — "
         "nos tests serveur simulent la base de données, donc ils ne peuvent pas prouver "
         "un comportement qui appartient à PostgreSQL. Le rendre automatique demanderait "
         "une vraie base dans la chaîne d'intégration : c'est la première ligne de nos "
         "pistes d'amélioration.\n\n"
         "Un choix que j'assume : le lint et le formatage sont rapportés mais ne bloquent "
         "pas. Nous avons de la dette là-dessus. Si nous les rendions bloquants "
         "aujourd'hui, la CI serait rouge en permanence — et une CI toujours rouge, plus "
         "personne ne la regarde. Nous préférons qu'elle reste crédible sur ce qui compte : "
         "la compilation et les tests. »")

# =========================================================================
# DIAPO 12 — Bogues et corrections
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 10, "Bogues rencontrés et corrigés",
              "Les quatre qui nous ont le plus appris")
carte(s, Inches(0.45), Inches(1.65), Inches(6.1), Inches(2.4),
      "Le 404 invisible en local",
      ["Symptôme : en ligne, tous les appels /api renvoyaient 404 —",
       "   alors que tout fonctionnait en local",
       "Cause : le proxy transmettait le mauvais en-tête Host.",
       "   Azure route selon cet en-tête ; Docker en local, non",
       "Correction : transmettre l'hôte de destination, pas celui",
       "   demandé par le navigateur"], ROUGE)
carte(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(2.4),
      "Le conteneur qui refusait de démarrer",
      ["Symptôme : l'image se construisait, le conteneur mourait",
       "   au lancement, sans message clair",
       "Cause : nos postes sont sous Windows ; le script de",
       "   démarrage était récupéré avec des fins de ligne CRLF",
       "Correction : forcer les fins de ligne LF sur les scripts",
       "   shell, au niveau du dépôt"], AMBRE)
carte(s, Inches(0.45), Inches(4.3), Inches(6.1), Inches(2.4),
      "L'application blanche au démarrage",
      ["Symptôme : erreur Angular NG0200 au chargement",
       "Cause : une dépendance circulaire dans la façade",
       "   d'authentification, déclenchée par le rafraîchissement",
       "   du profil au démarrage",
       "Correction : casser le cycle en différant la récupération",
       "   du profil"], BLEU_CLAIR)
carte(s, Inches(6.75), Inches(4.3), Inches(6.1), Inches(2.4),
      "Le bogue de méthode : la dérive d'intégration",
      ["Symptôme : cinq tickets « Terminé » dans Jira, absents",
       "   du produit",
       "Cause : des branches jamais fusionnées, jusqu'à 56",
       "   commits de retard",
       "Correction : réintégration en trois pull requests revues,",
       "   et Jira remis en cohérence avec le dépôt"], SARCELLE)
pied(s, "Larbi Saib", "1 min")
notes(s, "LARBI — 1 minute. En citer DEUX à l'oral, pas quatre.\n\n"
         "« Je vous en présente deux, les plus instructifs.\n\n"
         "Le premier nous a coûté une soirée. Une fois en ligne, tous les appels à l'API "
         "renvoyaient 404 — alors que tout marchait parfaitement en local. La cause : notre "
         "proxy transmettait le mauvais en-tête Host. Azure route selon cet en-tête, alors "
         "que Docker en local ne route pas comme ça du tout. La leçon : un environnement de "
         "développement qui fonctionne ne prouve rien sur la production.\n\n"
         "Le second n'est pas un bogue de code, c'est un bogue de méthode : la dérive "
         "d'intégration dont je parlais tout à l'heure. Cinq tickets marqués Terminé sans "
         "code dans le produit. Nous l'avons corrigé en réintégrant proprement et en "
         "remettant Jira en cohérence avec le dépôt. C'est celui qui nous a le plus "
         "appris. »")

# =========================================================================
# DIAPO 13 — Limites
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 11, "Nos limites", "Ce que le produit ne fait pas — et nous le savons")
carte(s, Inches(0.45), Inches(1.65), Inches(4.0), Inches(4.9),
      "Fonctionnalités absentes",
      ["Le décalage en bloc des rendez-",
       "   vous d'un médecin : le code",
       "   existe sur une branche, il n'est",
       "   pas intégré",
       "",
       "La configuration d'une clinique",
       "   par l'interface",
       "",
       "La gestion des médecins par",
       "   l'administrateur (CRUD)",
       "",
       "Aucune notification vers",
       "   l'extérieur : ni courriel, ni SMS"], AMBRE)
carte(s, Inches(4.65), Inches(1.65), Inches(4.0), Inches(4.9),
      "Limites techniques",
      ["Le multi-clinique est prévu dans le",
       "   modèle et appliqué côté sécurité,",
       "   mais peu outillé côté interface",
       "",
       "Pas de pagination serveur : le",
       "   volume académique le permet,",
       "   une vraie clinique non",
       "",
       "Dette de formatage : 233 fichiers",
       "   jamais passés au formateur,",
       "   9 avertissements d'accessibilité",
       "",
       "Démarrage à froid de quelques",
       "   secondes (scale-to-zero)"], BLEU_CLAIR)
carte(s, Inches(8.85), Inches(1.65), Inches(4.0), Inches(4.9),
      "Limites de méthode",
      ["Notre définition de « terminé » a",
       "   été trop souple pendant une",
       "   partie du projet — le code",
       "   pouvait exister sans être intégré",
       "",
       "Les revues de code sont arrivées",
       "   tard dans le projet",
       "",
       "Pas de tests de bout en bout",
       "   automatisés : le parcours complet",
       "   est validé à la main",
       "",
       "Aucun test de charge"], GRIS)
pied(s, "Zakaria Lahouiri", "45 s")
notes(s, "ZAKARIA — 45 secondes. Ton factuel, pas d'excuses.\n\n"
         "« Trois familles de limites, et je les assume.\n\n"
         "Des fonctionnalités manquent. La plus visible : le décalage en bloc des "
         "rendez-vous d'un médecin, pour absorber un retard. Le code existe sur une "
         "branche, mais il n'est pas intégré, et nous avons préféré ne pas le fusionner en "
         "catastrophe à trois jours de cette présentation plutôt que de risquer une "
         "régression sur le cœur du produit.\n\n"
         "Techniquement, nous n'avons pas de pagination serveur — ça tient à notre volume, "
         "pas à celui d'une vraie clinique. Et nous traînons de la dette de formatage.\n\n"
         "Et sur la méthode, notre définition de terminé a été trop souple pendant une "
         "partie du projet. C'est corrigé, mais ça nous a coûté du temps. »")

# =========================================================================
# DIAPO 14 — Pistes et strategie de progression
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 12, "La suite", "Par quoi nous continuerions, et dans quel ordre")
carte(s, Inches(0.45), Inches(1.65), Inches(3.9), Inches(2.4),
      "1.  Solder l'existant",
      ["Intégrer le décalage en bloc,",
       "   déjà écrit et testé",
       "Configuration de clinique et",
       "   gestion des médecins",
       "Passer le formateur sur tout",
       "   le dépôt, une fois pour toutes"], SARCELLE)
carte(s, Inches(4.55), Inches(1.65), Inches(3.9), Inches(2.4),
      "2.  Fiabiliser",
      ["Tests de bout en bout sur le",
       "   parcours complet",
       "Pagination serveur sur les",
       "   listes de rendez-vous",
       "Rendre le lint bloquant, une",
       "   fois la dette résorbée"], BLEU_CLAIR)
carte(s, Inches(8.65), Inches(1.65), Inches(4.2), Inches(2.4),
      "3.  Étendre la valeur",
      ["Rappels par courriel ou SMS —",
       "   la vraie arme anti-absence",
       "Liste d'attente : proposer",
       "   automatiquement un créneau",
       "   libéré par une annulation",
       "Portail patient en libre-service"], BLEU)
carte(s, Inches(0.45), Inches(4.3), Inches(12.4), Inches(2.35),
      "Notre stratégie de progression",
      ["Le produit est déjà déployable en une commande, et le schéma de base se reconstruit à l'identique. C'est le socle qui rend la suite possible.",
       "Nous garderions la même règle : une fonctionnalité par personne, menée de la base de données jusqu'à l'écran, fusionnée dans la semaine.",
       "L'ordre n'est pas arbitraire : d'abord finir ce qui est écrit, puis sécuriser ce qui existe, et seulement ensuite ajouter. Ajouter sur une base fragile, c'est ce qui nous a coûté le plus cher.",
       "La priorité métier suivante serait les rappels automatiques : c'est la fonctionnalité qui attaque directement le taux d'absence, le seul chiffre que la clinique voit sur sa facture."], BLEU)
pied(s, "Zakaria Lahouiri", "45 s")
notes(s, "ZAKARIA — 45 secondes.\n\n"
         "« Si nous continuions, nous procéderions en trois temps, et l'ordre compte.\n\n"
         "D'abord finir ce qui est déjà écrit — le décalage en bloc est codé et testé, il "
         "ne demande que d'être intégré. Ensuite fiabiliser : des tests de bout en bout, de "
         "la pagination. Et seulement après, ajouter de la valeur.\n\n"
         "Cet ordre, nous ne l'inventons pas : ajouter sur une base fragile, c'est "
         "exactement ce qui nous a coûté le plus cher ce semestre.\n\n"
         "Et si nous ne devions garder qu'une fonctionnalité, ce serait les rappels "
         "automatiques. C'est ce qui attaque le taux d'absence — le seul de nos chiffres "
         "que la clinique voit passer sur sa facture. »")

# =========================================================================
# DIAPO 15 — Contributions individuelles
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 13, "Qui a fait quoi",
              "Chacun présente sa contribution, ses difficultés et ce qu'il en retient")
carte(s, Inches(0.45), Inches(1.65), Inches(4.0), Inches(4.9),
      "Souleymane DIALLO",
      ["Réalisé : cahier des charges et",
       "   7 cas d'utilisation, socle du",
       "   monorepo et Docker,",
       "   authentification JWT, RBAC à",
       "   4 rôles, design system, passage",
       "   de la contrainte d'unicité à un",
       "   index PARTIEL (qui rend possible",
       "   l'annulation), déploiement Azure",
       "",
       "Difficulté : le 404 en production,",
       "   invisible en local",
       "",
       "Solution : comparer le routage",
       "   réel d'Azure et de Docker,",
       "   au lieu de chercher dans le code"], BLEU)
carte(s, Inches(4.65), Inches(1.65), Inches(4.0), Inches(4.9),
      "Zakaria Lahouiri",
      ["Réalisé : diagramme de classes,",
       "   disponibilités des médecins et",
       "   génération des créneaux, flux",
       "   clinique du jour, notifications",
       "   internes, export CSV",
       "",
       "Difficulté : mes branches avaient",
       "   pris jusqu'à 56 commits de",
       "   retard sur une interface refondue",
       "",
       "Solution : réintégration par pull",
       "   request, conflit par conflit, en",
       "   conservant les deux apports"], SARCELLE)
carte(s, Inches(8.85), Inches(1.65), Inches(4.0), Inches(4.9),
      "Larbi Saib",
      ["Réalisé : diagrammes de séquence,",
       "   modèle du patient léger, prise de",
       "   RDV par la réception, socle des",
       "   rendez-vous (entités, migration,",
       "   contrainte d'unicité), tableau de",
       "   bord et statistiques",
       "",
       "Difficulté : garantir qu'un créneau",
       "   ne soit jamais pris deux fois",
       "",
       "Solution : sortir la garantie du code",
       "   pour la poser en base, comme",
       "   contrainte d'unicité sur le créneau"], BLEU_CLAIR)
pied(s, "Chacun présente sa colonne", "1 min à trois")
notes(s, "LES TROIS — 20 secondes chacun. Chacun parle de SA colonne, à la première "
         "personne.\n\n"
         "Format imposé : ce que j'ai fait / ma difficulté / comment je l'ai résolue.\n\n"
         "SOULEYMANE : « J'ai porté le socle technique et la sécurité : le monorepo, "
         "l'authentification, le contrôle d'accès et la mise en ligne. Ma plus grosse "
         "difficulté a été un bogue qui n'existait qu'en production. J'ai arrêté de "
         "chercher dans le code et j'ai comparé comment Azure et Docker routent les "
         "requêtes — c'était là. »\n\n"
         "ZAKARIA : « J'ai porté les disponibilités, le flux du jour, les notifications et "
         "l'export. Ma difficulté a été de rattraper mes branches, qui avaient pris "
         "jusqu'à 56 commits de retard sur une interface entièrement refondue. Je les ai "
         "réintégrées par pull request, conflit par conflit, en gardant les deux apports "
         "plutôt qu'en écrasant. »\n\n"
         "LARBI : « J'ai porté le patient léger, la prise de rendez-vous par la réception et "
         "les statistiques. Ma difficulté a été de garantir qu'un créneau ne soit jamais "
         "pris deux fois. Vérifier dans le code ne suffit pas : entre la vérification et "
         "l'écriture, une autre requête peut passer. J'ai donc sorti la garantie du code "
         "pour la poser en base, comme contrainte d'unicité sur le créneau. Souleymane l'a "
         "ensuite transformée en index partiel, pour qu'une annulation libère le "
         "créneau. »\n\n"
         "CHACUN doit pouvoir répondre à une question sur SA colonne.")

# =========================================================================
# DIAPO 16 — Cloture
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(2.35), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(1.25), Inches(11.3), Inches(1.1))
para(tf, "MediPlan", taille=46, couleur=BLANC, gras=True, premier=True,
     aligne=PP_ALIGN.CENTER, espace_avant=0)
chiffre(s, Inches(1.4), Inches(2.85), Inches(2.5), "8", "Fonctionnalités livrées", SARCELLE)
chiffre(s, Inches(4.15), Inches(2.85), Inches(2.5), "186", "Tests verts", BLEU_CLAIR)
chiffre(s, Inches(6.9), Inches(2.85), Inches(2.5), "21", "Pull requests", BLEU_CLAIR)
chiffre(s, Inches(9.65), Inches(2.85), Inches(2.5), "0 $", "Coût mensuel", SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(4.75), Inches(11.3), Inches(1.6))
para(tf, "Une application en ligne, testée, reproductible en une commande.",
     taille=22, couleur=BLANC, gras=True, premier=True, aligne=PP_ALIGN.CENTER,
     espace_avant=10)
para(tf, "Souleymane DIALLO  ·  Zakaria Lahouiri  ·  Larbi Saib",
     taille=16, couleur=BLEU_PALE, aligne=PP_ALIGN.CENTER, espace_avant=14)
tb, tf = zone_texte(s, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.7))
para(tf, "Merci — nous répondons à vos questions", taille=20, couleur=BLEU_PALE2,
     premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
notes(s, "SOULEYMANE — 20 secondes, puis on ouvre les questions.\n\n"
         "« Pour résumer : huit fonctionnalités livrées et intégrées, 186 tests verts, une "
         "application en ligne pour environ zéro dollar par mois, et une infrastructure "
         "entièrement décrite en code — n'importe qui peut la redéployer d'une commande.\n\n"
         "Merci de votre attention. Nous répondons à vos questions. »\n\n"
         "RÈGLE POUR LES QUESTIONS : celui ou celle dont c'est le domaine répond. On ne se "
         "coupe pas la parole. Si on ne sait pas : « je ne l'ai pas vérifié, voici comment "
         "je le vérifierais » — jamais d'invention.\n\n"
         "QUESTIONS PROBABLES :\n"
         "· Doubles réservations ? → index unique partiel, PostgreSQL tranche ; vérifié en "
         "concurrence sur l'app en ligne, 201 + 409. Manuel, pas automatisé (Larbi)\n"
         "· Pourquoi le patient ne réserve pas ? → il appelle, c'est le flux réel (Souleymane)\n"
         "· Données de santé et sécurité ? → JWT, RBAC, chaque requête bornée à la clinique, "
         "mots de passe hachés, aucun secret dans le dépôt (Souleymane)\n"
         "· Pourquoi Azure et pas Railway ? → crédit étudiant non renouvelable, scale-to-zero, "
         "coût ~0 $ (Souleymane)\n"
         "· Le décalage en bloc ? → codé, non intégré, assumé (Zakaria)\n"
         "· Comment testez-vous ? → 186 tests, CI bloquante sur build et tests (Larbi)")

# =========================================================================
sortie = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "MediPlan-Presentation-Finale.pptx")
prs.save(sortie)
print(f"Genere : {sortie}")
print(f"Diapos : {len(prs.slides.__iter__.__self__._sldIdLst)}")
