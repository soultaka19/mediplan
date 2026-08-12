# -*- coding: utf-8 -*-
"""Genere le support de presentation du plan de Sprint 3 de MediPlan (PowerPoint).

Sprint 3 = le « dernier gros sprint » attendu par la consigne du cours.
Voir Sprint3-plan-de-sprint.md pour le detail et la note de renumerotation.

Usage : python docs/presentation/build_sprint3_pptx.py
Sortie : docs/presentation/MediPlan-Sprint3-Plan.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette (identique au support Sprint 1, pour la continuite visuelle) ---
BLEU = RGBColor(0x15, 0x4D, 0x8C)
BLEU_CLAIR = RGBColor(0x1E, 0x88, 0xE5)
SARCELLE = RGBColor(0x00, 0x89, 0x7B)
GRIS = RGBColor(0x40, 0x40, 0x40)
GRIS_CLAIR = RGBColor(0x6B, 0x6B, 0x6B)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
FOND = RGBColor(0xF4, 0xF7, 0xFB)
BLEU_PALE = RGBColor(0xCF, 0xE3, 0xFF)
BLEU_PALE2 = RGBColor(0x9F, 0xC8, 0xF5)
AMBRE = RGBColor(0xB5, 0x6A, 0x00)

LARGEUR = Inches(13.333)
HAUTEUR = Inches(7.5)

prs = Presentation()
prs.slide_width = LARGEUR
prs.slide_height = HAUTEUR
BLANK = prs.slide_layouts[6]


def fond(slide, couleur=BLANC):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = couleur


def rect(slide, x, y, w, h, couleur):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = couleur
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def zone_texte(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def para(tf, texte, taille=18, couleur=GRIS, gras=False, puce=False,
         espace_avant=6, niveau=0, aligne=PP_ALIGN.LEFT, premier=False):
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    p.alignment = aligne
    p.level = niveau
    p.space_after = Pt(espace_avant)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = (("•  " if puce else "") + texte)
    run.font.size = Pt(taille)
    run.font.bold = gras
    run.font.color.rgb = couleur
    run.font.name = "Calibri"
    return p


def notes(slide, texte):
    slide.notes_slide.notes_text_frame.text = texte


def bandeau_titre(slide, numero, titre, sous_titre=None):
    rect(slide, 0, 0, LARGEUR, Inches(1.25), BLEU)
    rect(slide, 0, Inches(1.25), LARGEUR, Pt(6), SARCELLE)
    past = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.32),
                                  Inches(0.62), Inches(0.62))
    past.fill.solid(); past.fill.fore_color.rgb = SARCELLE
    past.line.fill.background(); past.shadow.inherit = False
    pf = past.text_frame; pf.word_wrap = False
    pp = pf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); r.text = str(numero); r.font.size = Pt(22)
    r.font.bold = True; r.font.color.rgb = BLANC
    tb, tf = zone_texte(slide, Inches(1.25), Inches(0.18), Inches(11.6), Inches(1.0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, titre, taille=28, couleur=BLANC, gras=True, premier=True, espace_avant=0)
    if sous_titre:
        para(tf, sous_titre, taille=14, couleur=BLEU_PALE, espace_avant=0)


def pied(slide, orateur, duree):
    tb, tf = zone_texte(slide, Inches(0.45), Inches(7.0), Inches(12.4), Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = f"Orateur : {orateur}"
    r.font.size = Pt(11); r.font.color.rgb = GRIS_CLAIR; r.font.italic = True
    tb2, tf2 = zone_texte(slide, Inches(10.0), Inches(7.0), Inches(2.9), Inches(0.4))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"⏱  {duree}"
    r2.font.size = Pt(11); r2.font.color.rgb = SARCELLE; r2.font.bold = True


def carte(slide, x, y, w, h, titre, lignes, couleur_titre=BLEU):
    rect(slide, x, y, w, h, FOND)
    rect(slide, x, y, Pt(6), h, couleur_titre)
    tb, tf = zone_texte(slide, x + Inches(0.25), y + Inches(0.12),
                        w - Inches(0.4), h - Inches(0.2))
    para(tf, titre, taille=16, couleur=couleur_titre, gras=True, premier=True,
         espace_avant=4)
    for l in lignes:
        para(tf, l, taille=13, couleur=GRIS, puce=bool(l.strip()), espace_avant=3)


def chiffre(slide, x, y, largeur, valeur, libelle, couleur):
    rect(slide, x, y, largeur, Inches(1.5), couleur)
    tb, tf = zone_texte(slide, x, y + Inches(0.08), largeur, Inches(1.4))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, valeur, taille=40, couleur=BLANC, gras=True, premier=True,
         aligne=PP_ALIGN.CENTER, espace_avant=0)
    para(tf, libelle, taille=13, couleur=BLANC, aligne=PP_ALIGN.CENTER, espace_avant=0)


# =========================================================================
# DIAPO 1 — Titre
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(3.05), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(1.4))
para(tf, "MediPlan", taille=60, couleur=BLANC, gras=True, premier=True,
     aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(0.7))
para(tf, "Plateforme web de gestion des rendez-vous médicaux", taille=22,
     couleur=BLEU_PALE, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(3.3), Inches(11.3), Inches(1.4))
para(tf, "Plan de Sprint 3", taille=28, couleur=BLANC,
     gras=True, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=8)
para(tf, "Le cycle de vie complet du rendez-vous — notre dernier gros sprint", taille=18,
     couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.4))
para(tf, "Équipe : Souleymane DIALLO  ·  Zakaria Lahouiri  ·  Larbi Saib", taille=18,
     couleur=BLANC, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=4)
para(tf, "Projet intégrateur — Programmation informatique  ·  Collège La Cité  ·  Printemps 2026",
     taille=14, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=2)
notes(s, "Bonjour. Nous venons vous présenter notre plan de Sprint 3. Nous allons d'abord "
         "prendre deux minutes pour vous montrer où le projet en est concrètement, parce "
         "que ça explique le contenu du sprint que nous proposons. Ensuite nous verrons "
         "l'objectif, le backlog, la répartition et les risques. Environ dix minutes.")

# =========================================================================
# DIAPO 2 — Ou en est le projet
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 1, "Où en est le projet",
              "Sprint 2 terminé — livré et intégré, pas seulement codé")
chiffre(s, Inches(0.45), Inches(1.6), Inches(2.85), "4", "Fonctionnalités livrées", SARCELLE)
chiffre(s, Inches(3.55), Inches(1.6), Inches(2.85), "171", "Tests automatisés verts", BLEU)
chiffre(s, Inches(6.65), Inches(1.6), Inches(2.85), "2", "Sprints terminés", BLEU_CLAIR)
chiffre(s, Inches(9.75), Inches(1.6), Inches(2.85), "6", "Tickets au Sprint 3", GRIS)
carte(s, Inches(0.45), Inches(3.35), Inches(6.1), Inches(3.4),
      "Livré et fusionné dans dev (et promu sur main)",
      ["Authentification complète — MEDIPLAN-15, 16, 17",
       "   inscription, connexion JWT, verrouillage de compte,",
       "   réinitialisation du mot de passe, RBAC à 4 rôles",
       "Patient léger, géré par la réception — MEDIPLAN-35",
       "Disponibilités médecins + génération des créneaux — MEDIPLAN-20",
       "Prise de RDV par la réception + flux du jour — MEDIPLAN-21, 36, 23",
       "Socle technique RDV + index anti-double-booking — MEDIPLAN-49"], SARCELLE)
carte(s, Inches(6.75), Inches(3.35), Inches(6.1), Inches(3.4),
      "Nos preuves (montrées à l'écran)",
      ["Jira : tableaux Sprint 1 et Sprint 2, tickets en « Terminé »",
       "GitHub : branche dev = main, merges 71cb40f puis 38400da",
       "Tests : 48 backend + 123 frontend, tous verts",
       "Démo live : connexion → agenda → réservation → flux du jour",
       " ",
       "Le cœur métier anti-double-réservation (OM-04) est garanti",
       "   par un index unique partiel en base, testé en concurrence"], BLEU)
pied(s, "Souleymane", "2 min")
notes(s, "Commençons par l'état réel. Notre Sprint 2 est terminé, et il contient plus que "
         "l'authentification : nous avons quatre fonctionnalités livrées et fusionnées dans "
         "notre branche dev — pas juste codées sur un poste. L'authentification complète "
         "avec JWT, verrouillage de compte et RBAC ; le patient léger géré par la "
         "réception ; les disponibilités des médecins avec génération automatique des "
         "créneaux ; et la prise de rendez-vous par la réception avec le flux du jour. "
         "Nous pouvons vous le démontrer en direct : on se connecte, on ouvre l'agenda, on "
         "réserve un créneau, et il apparaît dans le flux du jour. "
         "Le cœur métier, l'anti-double-réservation, est garanti par un index unique en "
         "base et testé en concurrence. Nous avons 171 tests automatisés verts, et tout est "
         "promu sur main.")

# =========================================================================
# DIAPO 3 — Objectif & backlog
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 2, "Objectif et backlog du Sprint 3",
              "Priorisation MoSCoW — notre dernier gros sprint")
rect(s, Inches(0.45), Inches(1.6), Inches(12.4), Inches(1.1), SARCELLE)
tb, tf = zone_texte(s, Inches(0.75), Inches(1.72), Inches(11.8), Inches(0.95))
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "Objectif du sprint", taille=12, couleur=BLEU_PALE, gras=True, premier=True,
     espace_avant=2)
para(tf, "La réception gère le cycle de vie complet d'un rendez-vous — création, annulation "
         "motivée, notification à l'équipe — et la clinique lit son activité réelle.",
     taille=16, couleur=BLANC, gras=True, espace_avant=0)
carte(s, Inches(0.45), Inches(2.95), Inches(8.5), Inches(3.8),
      "Backlog priorisé",
      ["MUST · MEDIPLAN-22 — Annuler un RDV avec motif obligatoire, libère le créneau",
       "MUST · MEDIPLAN-25 — Notifications internes sur les changements de RDV",
       "MUST · MEDIPLAN-26 — Statistiques réelles (remplacent les compteurs provisoires)",
       "MUST · MEDIPLAN-50 — Sécurité RDV : patient léger non authentifiable,",
       "        cloisonnement des cliniques",
       "SHOULD · MEDIPLAN-24 — Décaler en bloc les RDV (code écrit, reste à intégrer)",
       "SHOULD · MEDIPLAN-51 — Interface modernisée (palette clinique, typographie)",
       " ",
       "Un Must par personne. Les Should servent de variable d'ajustement."], BLEU)
carte(s, Inches(9.15), Inches(2.95), Inches(3.7), Inches(3.8),
      "Hors périmètre — assumé",
      ["Export CSV (27)",
       "CI GitHub Actions (30)",
       "Déploiement Railway (40)",
       "Documentation des tests",
       "Réflexion UI/UX",
       " ",
       "→ tout cela au Sprint 4",
       "   de finalisation"], GRIS)
pied(s, "Souleymane", "2 min")
notes(s, "Voici l'objectif : à la fin du sprint, la réception gère le cycle de vie complet "
         "d'un rendez-vous — le créer, l'annuler avec un motif, et l'équipe en est notifiée "
         "— et la clinique voit son activité réelle sur son tableau de bord. C'est la suite "
         "logique de ce que nous avons livré : aujourd'hui on sait créer un rendez-vous, "
         "mais pas encore l'annuler proprement. "
         "Quatre tickets en Must. L'annulation avec motif, qui libère le créneau. Les "
         "notifications internes. Les statistiques réelles, qui remplacent les compteurs "
         "provisoires. Et un ticket de sécurité : garantir qu'un patient léger ne peut "
         "jamais s'authentifier, et qu'on n'accède pas aux rendez-vous d'une autre clinique "
         "— nous manipulons des données de santé, nous ne voulions pas le repousser. "
         "Deux Should : le décalage en bloc, dont le code est écrit mais pas encore "
         "intégré, et l'interface modernisée. "
         "Nous assumons de sortir le reste : export CSV, intégration continue, déploiement "
         "et documentation partent au sprint de finalisation. Nous préférons un sprint tenu "
         "qu'un sprint ambitieux à moitié fait.")

# =========================================================================
# DIAPO 4 — Repartition & sequencement
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 3, "Répartition et séquençage", "Un engagement ferme par membre")
carte(s, Inches(0.45), Inches(1.7), Inches(4.0), Inches(2.6),
      "Souleymane DIALLO",
      ["MEDIPLAN-22 — annulation (Must)",
       "MEDIPLAN-50 — sécurité RDV (Must)",
       "MEDIPLAN-51 — interface (Should)",
       " ",
       "Rôle : pilotage, intégration"], BLEU)
carte(s, Inches(4.65), Inches(1.7), Inches(4.0), Inches(2.6),
      "Zakaria Lahouiri",
      ["MEDIPLAN-25 — notifications",
       "   internes (Must) — En cours",
       "MEDIPLAN-24 — décalage en bloc",
       "   (Should) — à intégrer"], SARCELLE)
carte(s, Inches(8.85), Inches(1.7), Inches(4.0), Inches(2.6),
      "Larbi Saib",
      ["MEDIPLAN-26 — statistiques et",
       "   KPI réels (Must)",
       " ",
       "Retire les compteurs provisoires",
       "du tableau de bord"], BLEU_CLAIR)
carte(s, Inches(0.45), Inches(4.5), Inches(7.4), Inches(2.25),
      "Séquençage — chaque palier est démontrable",
      ["PR-1  ·  MEDIPLAN-22 passe SEUL et EN PREMIER (il touche le statut du RDV)",
       "PR-2  ·  MEDIPLAN-25 — dépend de l'événement d'annulation",
       "PR-3  ·  MEDIPLAN-26 — en parallèle (lecture seule, pas de conflit)",
       "PR-4  ·  MEDIPLAN-50 — revue de sécurité transverse, en fin de sprint",
       "PR-UI ·  MEDIPLAN-51 — fusionnée tôt, pour éviter les conflits de styles"], BLEU)
carte(s, Inches(8.05), Inches(4.5), Inches(4.8), Inches(2.25),
      "Dépendance bloquante identifiée",
      ["L'endpoint de changement de statut",
       "d'un RDV doit d'abord accepter la",
       "valeur « cancelled ».",
       " ",
       "→ C'est le tout premier travail du",
       "   sprint, sinon MEDIPLAN-22 et 25",
       "   sont bloqués tous les deux."], AMBRE)
pied(s, "Zakaria", "2 min")
notes(s, "Côté répartition, chaque membre porte un engagement ferme, un Must, démontrable "
         "indépendamment. Souleymane prend l'annulation et le ticket de sécurité, plus "
         "l'interface en Should. Je prends les notifications internes — c'est déjà en cours "
         "— et la reprise du décalage en bloc. Larbi prend les statistiques réelles. "
         "Le séquençage vient d'une leçon du sprint précédent. Le ticket 22 touche le "
         "statut du rendez-vous : il passe seul et en premier. Ensuite les notifications et "
         "les statistiques avancent en parallèle, parce qu'ils ne se marchent pas dessus. "
         "La sécurité se fait en revue transverse à la fin. "
         "Nous avons identifié une dépendance bloquante : l'endpoint de statut doit d'abord "
         "accepter la valeur « annulé ». Sans ça, deux tickets sont bloqués. C'est donc le "
         "tout premier travail du sprint.")

# =========================================================================
# DIAPO 5 — Risques & regles de travail
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 4, "Risques et règles de travail",
              "Ce que le sprint précédent nous a appris")
carte(s, Inches(0.45), Inches(1.7), Inches(6.1), Inches(2.5),
      "⚠  Risque n° 1 — la dérive d'intégration (vécue)",
      ["Six branches ont vécu en parallèle sans être fusionnées",
       "Résultat : deux migrations sur le même horodatage, et le",
       "   module de rendez-vous recréé deux fois par deux personnes",
       "Coût : une réintégration manuelle, branche par branche",
       "Séquelle encore visible : MEDIPLAN-24, dont le code existe",
       "   mais n'est pas intégré → rouvert et repris dans ce sprint"], AMBRE)
carte(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(2.5),
      "⚠  Risque n° 2 — rythme inégal entre membres",
      ["Mitigation : un Must par personne, chacun démontrable seul",
       "Si un ticket décroche, l'objectif reste atteignable en",
       "   dégradant les Should (interface, décalage en bloc)",
       "Point d'intégration à mi-sprint pour détecter tôt",
       "Aucun ticket ne dépend de deux personnes à la fois"], BLEU_CLAIR)
carte(s, Inches(0.45), Inches(4.45), Inches(12.4), Inches(2.3),
      "Nos 6 règles pour le Sprint 3",
      ["1. Socle d'abord — le ticket qui touche le modèle ou une migration passe seul et en premier",
       "2. PR sous 72 h — aucune branche ne vit plus de trois jours sans PR ouverte vers dev",
       "3. Horodatage de migration annoncé à l'équipe avant création — séquence continue, jamais de doublon",
       "4. Un module par ticket — interdiction de recréer un module qui existe déjà",
       "5. Gates de fusion — lint, tests verts, et migrations exécutées sur une base neuve",
       "6. Point d'intégration à mi-sprint — pas à la fin"], SARCELLE)
pied(s, "Larbi", "2 min")
notes(s, "Nous voulons être transparents sur un point : notre sprint précédent a mal "
         "commencé. Six branches ont vécu en parallèle sans jamais être fusionnées. "
         "Résultat : deux migrations créées sur le même horodatage, et le module de "
         "rendez-vous développé deux fois par deux personnes différentes. Il a fallu tout "
         "réintégrer à la main. C'est réglé, tout est dans dev et les tests sont verts, "
         "mais cela nous a coûté du temps. Il reste même une séquelle : le ticket 24, le "
         "décalage en bloc — le code est écrit, mais il n'est toujours pas intégré. En "
         "préparant ce sprint, nous l'avons rouvert honnêtement plutôt que de le laisser "
         "marqué terminé à tort. "
         "Nous en avons tiré six règles, appliquées dès ce sprint : le ticket qui touche la "
         "base passe seul et en premier ; aucune branche ne vit plus de trois jours sans "
         "PR ; on annonce les horodatages de migration ; un seul module par ticket ; on ne "
         "fusionne que si le lint, les tests et les migrations passent ; et on fait un "
         "point d'intégration au milieu du sprint, pas à la fin.")

# =========================================================================
# DIAPO 6 — Definition of Done & suite
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 5, "« Terminé » veut dire quoi, et après ?",
              "Definition of Done · démo visée · Sprint 4")
carte(s, Inches(0.45), Inches(1.7), Inches(6.1), Inches(2.9),
      "Un ticket est « Terminé » quand…",
      ["le code est FUSIONNÉ dans dev — pas seulement poussé",
       "le lint et les tests automatisés sont verts",
       "les migrations s'exécutent sur une base neuve",
       "les critères d'acceptation Given/When/Then du ticket",
       "   Jira sont vérifiés un par un",
       "la fonctionnalité est démontrable dans l'application"], BLEU)
carte(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(2.9),
      "La démo que nous viserons en fin de sprint",
      ["1. Réserver un rendez-vous pour un patient",
       "2. L'annuler en saisissant un motif obligatoire",
       "3. Voir le créneau se libérer et redevenir réservable",
       "4. Voir la notification apparaître pour l'équipe",
       "5. Voir le compteur du tableau de bord se mettre à jour",
       " ",
       "Un scénario, cinq fonctionnalités enchaînées."], SARCELLE)
carte(s, Inches(0.45), Inches(4.85), Inches(12.4), Inches(1.9),
      "Sprint 4 — finalisation",
      ["Export CSV des rendez-vous (MEDIPLAN-27)  ·  CI GitHub Actions : lint, tests, build (MEDIPLAN-30)",
       "Déploiement Railway (MEDIPLAN-40, bonus)  ·  Documentation des tests  ·  Réflexion UI/UX d'équipe",
       "Préparation de la démonstration finale"], BLEU_CLAIR)
pied(s, "Larbi", "1 min 30")
notes(s, "Nous avons voulu être précis sur ce que « terminé » veut dire chez nous, parce "
         "que c'est exactement là que nous nous sommes fait piéger. Un ticket est terminé "
         "quand le code est fusionné dans dev — pas juste poussé sur une branche — quand le "
         "lint et les tests sont verts, quand les migrations passent sur une base neuve, "
         "quand les critères d'acceptation Jira sont vérifiés un par un, et quand la "
         "fonctionnalité est démontrable. "
         "La démo que nous viserons tient en un scénario : on réserve un rendez-vous, on "
         "l'annule avec un motif, on voit le créneau se libérer, la notification "
         "apparaître, et le compteur du tableau de bord bouger. Un seul scénario qui "
         "enchaîne cinq fonctionnalités. "
         "Ensuite, le Sprint 4 sera consacré à la finalisation : export CSV, intégration "
         "continue, déploiement, documentation des tests et réflexion UI/UX.")

# =========================================================================
# DIAPO 7 — Conclusion / Questions
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(3.5), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.4))
para(tf, "Merci de votre attention", taille=40, couleur=BLANC, gras=True,
     premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(3.8), Inches(11.3), Inches(1.6))
para(tf, "Questions ?", taille=28, couleur=BLEU_PALE, gras=True,
     premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
para(tf, "Équipe MediPlan — Souleymane DIALLO  ·  Zakaria Lahouiri  ·  Larbi Saib",
     taille=16, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=14)
notes(s, "En résumé : notre Sprint 2 est terminé avec quatre fonctionnalités livrées et "
         "intégrées ; le Sprint 3 a un objectif clair, le cycle de vie complet du "
         "rendez-vous, quatre tickets Must, un par personne, et des règles de travail "
         "tirées de nos erreurs. Nous sommes disponibles pour vos questions."
         "\n\nQuestions probables et réponses :\n"
         "— « Est-ce réaliste ? » → 4 Must, un par personne, aucun ticket partagé, et deux "
         "Should qui servent de variable d'ajustement.\n"
         "— « Quelles sont les dates du sprint ? » → À caler avec vous aujourd'hui, en "
         "visant la fin des cours pour garder le Sprint 4 de finalisation.\n"
         "— « Pourquoi le patient ne réserve-t-il pas lui-même ? » → Volontairement "
         "repoussé : nous avons priorisé la réception, qui est l'utilisateur réel du "
         "cahier des charges (le patient appelle au téléphone).\n"
         "— « Le ticket 24 était-il vraiment terminé ? » → Non, il était marqué terminé à "
         "tort : le code existe mais n'était pas intégré. Nous l'avons rouvert nous-mêmes "
         "en préparant ce sprint.")

# --- Sauvegarde ----------------------------------------------------------
sortie = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "MediPlan-Sprint3-Plan.pptx")
prs.save(sortie)
print("Genere :", sortie)
print("Nombre de diapos :", len(prs.slides._sldIdLst))
