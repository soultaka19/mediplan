# -*- coding: utf-8 -*-
"""Genere le support de presentation Sprint 1 de MediPlan (PowerPoint).

Usage : python docs/presentation/build_pptx.py
Sortie : docs/presentation/MediPlan-Sprint1-Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette (bleu medical + accent sarcelle) -----------------------------
BLEU = RGBColor(0x15, 0x4D, 0x8C)
BLEU_CLAIR = RGBColor(0x1E, 0x88, 0xE5)
SARCELLE = RGBColor(0x00, 0x89, 0x7B)
GRIS = RGBColor(0x40, 0x40, 0x40)
GRIS_CLAIR = RGBColor(0x6B, 0x6B, 0x6B)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
FOND = RGBColor(0xF4, 0xF7, 0xFB)
BLEU_PALE = RGBColor(0xCF, 0xE3, 0xFF)
BLEU_PALE2 = RGBColor(0x9F, 0xC8, 0xF5)

LARGEUR = Inches(13.333)
HAUTEUR = Inches(7.5)

prs = Presentation()
prs.slide_width = LARGEUR
prs.slide_height = HAUTEUR
BLANK = prs.slide_layouts[6]


def fond(slide, couleur=BLANC):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = couleur


def rect(slide, x, y, w, h, couleur):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = couleur
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def zone_texte(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def para(tf, texte, taille=18, couleur=GRIS, gras=False, puce=False,
         espace_avant=6, niveau=0, aligne=PP_ALIGN.LEFT, premier=False):
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    p.alignment = aligne
    p.level = niveau
    p.space_after = Pt(espace_avant)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = (("•  " if puce else "") + texte)
    run.font.size = Pt(taille)
    run.font.bold = gras
    run.font.color.rgb = couleur
    run.font.name = "Calibri"
    return p


def notes(slide, texte):
    slide.notes_slide.notes_text_frame.text = texte


def bandeau_titre(slide, numero, titre, sous_titre=None):
    rect(slide, 0, 0, LARGEUR, Inches(1.25), BLEU)
    rect(slide, 0, Inches(1.25), LARGEUR, Pt(6), SARCELLE)
    past = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.32),
                                  Inches(0.62), Inches(0.62))
    past.fill.solid(); past.fill.fore_color.rgb = SARCELLE
    past.line.fill.background(); past.shadow.inherit = False
    pf = past.text_frame; pf.word_wrap = False
    pp = pf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run(); r.text = str(numero); r.font.size = Pt(22)
    r.font.bold = True; r.font.color.rgb = BLANC
    tb, tf = zone_texte(slide, Inches(1.25), Inches(0.18), Inches(11.6), Inches(1.0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, titre, taille=28, couleur=BLANC, gras=True, premier=True, espace_avant=0)
    if sous_titre:
        para(tf, sous_titre, taille=14, couleur=BLEU_PALE, espace_avant=0)


def pied(slide, orateur, duree):
    tb, tf = zone_texte(slide, Inches(0.45), Inches(7.0), Inches(12.4), Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = f"Orateur : {orateur}"
    r.font.size = Pt(11); r.font.color.rgb = GRIS_CLAIR; r.font.italic = True
    tb2, tf2 = zone_texte(slide, Inches(10.0), Inches(7.0), Inches(2.9), Inches(0.4))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"⏱  {duree}"
    r2.font.size = Pt(11); r2.font.color.rgb = SARCELLE; r2.font.bold = True


def carte(slide, x, y, w, h, titre, lignes, couleur_titre=BLEU):
    rect(slide, x, y, w, h, FOND)
    rect(slide, x, y, Pt(6), h, couleur_titre)
    tb, tf = zone_texte(slide, x + Inches(0.25), y + Inches(0.12),
                        w - Inches(0.4), h - Inches(0.2))
    para(tf, titre, taille=16, couleur=couleur_titre, gras=True, premier=True,
         espace_avant=4)
    for l in lignes:
        para(tf, l, taille=13, couleur=GRIS, puce=bool(l.strip()), espace_avant=3)


# =========================================================================
# DIAPO 1 — Titre
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(3.05), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(1.4))
para(tf, "MediPlan", taille=60, couleur=BLANC, gras=True, premier=True,
     aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(0.7))
para(tf, "Plateforme web de gestion des rendez-vous médicaux", taille=22,
     couleur=BLEU_PALE, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(3.3), Inches(11.3), Inches(1.4))
para(tf, "Présentation d'avancement — Sprint 1", taille=28, couleur=BLANC,
     gras=True, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=8)
para(tf, "Analyse & conception", taille=18, couleur=BLEU_PALE2,
     aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.4))
para(tf, "Équipe : Souleymane DIALLO  ·  Zakaria Lahouiri  ·  Larbi Saib", taille=18,
     couleur=BLANC, premier=True, aligne=PP_ALIGN.CENTER, espace_avant=4)
para(tf, "Projet intégrateur — Programmation informatique  ·  Collège La Cité  ·  Printemps 2026",
     taille=14, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=2)
para(tf, "Mercredi 10 juin 2026", taille=14, couleur=BLEU_PALE2,
     aligne=PP_ALIGN.CENTER, espace_avant=2)
notes(s, "Bonjour, nous sommes Souleymane et Zakaria, et nous présentons l'avancement "
         "du Sprint 1 de notre projet MediPlan. Ce premier sprint portait sur l'analyse "
         "et la conception. Nous allons vous montrer en environ 10 minutes où nous en "
         "sommes, comment nous nous sommes organisés avec Jira et GitHub, et ce qu'il "
         "reste à faire. (Souleymane ouvre les diapos 1 à 6 ; Zakaria et Larbi prennent "
         "la parole à partir de la diapo 7 : Zakaria sur les classes, Larbi sur la séquence.)")

# =========================================================================
# DIAPO 2 — Rappel du projet
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 1, "Rappel du projet", "Titre · problématique · solution · utilisateurs")
carte(s, Inches(0.45), Inches(1.7), Inches(6.1), Inches(2.4),
      "Problématique",
      ["Cliniques : prise de RDV par téléphone, agendas papier, fichiers Excel",
       "Erreurs, doubles réservations, oublis, pas de vue d'ensemble",
       "Aucune disponibilité 24 h/24 pour les patients"], BLEU)
carte(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(2.4),
      "Solution proposée",
      ["Plateforme web unique, sécurisée, accessible 24 h/24",
       "Réservation / modification / annulation en ligne",
       "Gestion des disponibilités, flux du jour, statistiques",
       "Angular + NestJS + PostgreSQL + Docker"], SARCELLE)
carte(s, Inches(0.45), Inches(4.3), Inches(12.4), Inches(2.45),
      "Utilisateurs / clients visés  —  4 rôles (RBAC)",
      ["Patient : réserve, modifie ou annule ses rendez-vous sans téléphoner",
       "Médecin : consulte son horaire, gère ses disponibilités, suit le flux du jour",
       "Administrateur de clinique : gère utilisateurs, médecins, disponibilités, statistiques",
       "Super administrateur : configure les cliniques et la plateforme"], BLEU_CLAIR)
pied(s, "Souleymane", "1 min 30")
notes(s, "MediPlan répond à un problème concret des petites et moyennes cliniques : "
         "aujourd'hui la prise de rendez-vous se fait par téléphone, sur agendas papier "
         "ou fichiers Excel, ce qui cause des erreurs, des doubles réservations et aucune "
         "vue d'ensemble. Notre solution est une plateforme web unique et sécurisée, "
         "accessible 24 h/24, où le patient réserve lui-même. Quatre rôles interagissent "
         "avec le système selon un modèle de permissions RBAC : patient, médecin, "
         "administrateur de clinique et super administrateur.")

# =========================================================================
# DIAPO 3 — Objectif du Sprint 1
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 2, "Objectif du Sprint 1", "Analyse & conception")
carte(s, Inches(0.45), Inches(1.7), Inches(6.1), Inches(4.6),
      "Ce que l'équipe voulait accomplir",
      ["Cadrer le périmètre : cahier des charges (fonctions EF-01 à EF-09)",
       "Modéliser le système avant de coder",
       "Produire les diagrammes UML : cas d'utilisation, classes, séquence",
       "Ajouter un diagramme entité-association (ERD) pour PostgreSQL",
       "Mettre en place les outils de suivi : Jira + GitHub",
       "Déposer le dossier de conception et le soumettre sur eCité"], BLEU)
carte(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(4.6),
      "Pourquoi cet objectif est important",
      ["La conception est la fondation de tout le projet",
       "Éviter de coder dans le flou : moins d'erreurs ensuite",
       "Le diagramme de classes guide directement la base de données",
       "Identifier tôt les risques techniques (gestion des disponibilités)",
       "Aligner toute l'équipe sur une vision commune",
       "Découper le travail en Epics / sprints réalistes"], SARCELLE)
pied(s, "Souleymane", "1 min")
notes(s, "L'objectif du Sprint 1 était l'analyse et la conception : avant d'écrire la "
         "moindre ligne de code, nous voulions bien cadrer le périmètre à partir du "
         "cahier des charges, puis modéliser le système avec des diagrammes UML, et "
         "mettre en place nos outils de suivi Jira et GitHub. C'est important parce que "
         "la conception est la fondation du projet : elle nous évite de coder dans le "
         "flou, elle guide la conception de la base de données, et elle nous permet "
         "d'identifier tôt les parties risquées, comme la gestion des disponibilités.")

# =========================================================================
# DIAPO 4 — Organisation & demarche
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 3, "Notre organisation & démarche", "Comment l'équipe travaille")
carte(s, Inches(0.45), Inches(1.7), Inches(4.0), Inches(2.3),
      "Planification",
      ["Jira : Epics, user stories, tâches",
       "Sprints gérés par labels Sprint-1…6",
       "Statuts et responsables à jour"], BLEU)
carte(s, Inches(4.65), Inches(1.7), Inches(4.0), Inches(2.3),
      "Versionnement",
      ["GitHub : dépôt privé partagé",
       "Dossier docs/conception/",
       "Accès donné à la professeure"], SARCELLE)
carte(s, Inches(8.85), Inches(1.7), Inches(4.0), Inches(2.3),
      "Modélisation",
      ["Diagrammes en Mermaid (versionnables)",
       "Affichage direct dans GitHub",
       "Images PNG exportées"], BLEU_CLAIR)
carte(s, Inches(0.45), Inches(4.25), Inches(12.4), Inches(2.5),
      "Répartition des rôles dans l'équipe",
      ["Souleymane DIALLO : pilotage du projet, mise en place Jira & GitHub, cahier des charges, 7 diagrammes de cas d'utilisation",
       "Zakaria Lahouiri : diagramme de classes et son explication écrite",
       "Larbi Saib : diagramme de séquence et son explication écrite",
       "Travail en équipe : relecture croisée, validation des diagrammes et rédaction des explications"], BLEU)
pied(s, "Souleymane", "1 min")
notes(s, "Un mot sur notre démarche, car ce sprint, c'est surtout de l'organisation. "
         "Nous planifions tout dans Jira avec des Epics, des user stories et des tâches, "
         "et nous gérons nos sprints par des labels. Tout le code et les documents sont "
         "versionnés sur GitHub, dans le dossier docs/conception. Nos diagrammes sont "
         "écrits en Mermaid, ce qui permet de les versionner comme du texte et de les "
         "afficher directement dans GitHub. Côté répartition : je me suis occupé du "
         "pilotage, de Jira, GitHub et des cas d'utilisation ; Zakaria a pris en charge "
         "le diagramme de classes et Larbi le diagramme de séquence. Nous nous relisons "
         "mutuellement.")

# =========================================================================
# DIAPO 5 — Avancement dans Jira
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 4, "Avancement dans Jira", "Epics · user stories · tâches · progression")
def chiffre(slide, x, valeur, libelle, couleur):
    rect(slide, x, Inches(1.7), Inches(2.85), Inches(1.5), couleur)
    tb, tf = zone_texte(slide, x, Inches(1.78), Inches(2.85), Inches(1.4))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, valeur, taille=40, couleur=BLANC, gras=True, premier=True,
         aligne=PP_ALIGN.CENTER, espace_avant=0)
    para(tf, libelle, taille=13, couleur=BLANC, aligne=PP_ALIGN.CENTER, espace_avant=0)
chiffre(s, Inches(0.45), "7", "Epics", BLEU)
chiffre(s, Inches(3.55), "23", "User stories", BLEU_CLAIR)
chiffre(s, Inches(6.65), "3", "Tâches", SARCELLE)
chiffre(s, Inches(9.75), "6", "Sprints planifiés", GRIS)
carte(s, Inches(0.45), Inches(3.45), Inches(6.1), Inches(3.3),
      "Sprint 1 — Analyse & conception",
      ["Epic E1 (Analyse / Conception) : conception terminée",
       "Stories de conception : cahier des charges, cas d'usage,",
       "   classes, séquence, ERD, dépôt GitHub — terminées",
       "Responsables assignés pour chaque ticket",
       "Sprint 1 quasi complet, jalon conception atteint"], SARCELLE)
carte(s, Inches(6.75), Inches(3.45), Inches(6.1), Inches(3.3),
      "Backlog des sprints suivants (planifié)",
      ["S2 — Fondations & authentification (Docker, JWT, RBAC)",
       "S3 — Cliniques, médecins & disponibilités",
       "S4 — Rendez-vous & notifications",
       "S5 — Flux clinique & tableaux de bord",
       "S6 — Qualité, CI/CD & déploiement"], BLEU)
pied(s, "Souleymane", "1 min 30")
notes(s, "Voici notre plan de travail dans Jira. Nous avons structuré le projet en 7 "
         "Epics, 23 user stories et 3 tâches, répartis sur 6 sprints. Pour le Sprint 1, "
         "l'Epic Analyse et Conception est terminée : toutes les stories liées au cahier "
         "des charges, aux diagrammes et au dépôt GitHub sont passées en terminé, avec un "
         "responsable assigné sur chaque ticket. Le reste du backlog est déjà planifié "
         "jusqu'au sprint 6, ce qui nous donne une feuille de route claire. Notre tableau "
         "reflète l'état réel : les tâches ne sont pas toutes restées dans la colonne "
         "« À faire ».")

# =========================================================================
# DIAPO 6 — Conception (1) cas d'utilisation
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 5, "Dossier de conception (1/2)", "Diagrammes de cas d'utilisation")
tb, tf = zone_texte(s, Inches(0.45), Inches(1.5), Inches(12.4), Inches(0.6))
para(tf, "7 diagrammes produits (consigne : minimum 5) — couvrent les fonctions EF-01 à EF-09",
     taille=15, couleur=GRIS, gras=True, premier=True, espace_avant=0)
gauche = ["UC-1  Vue d'ensemble du système (4 acteurs)",
          "UC-2  Authentification & gestion des comptes",
          "UC-3  Réservation / modification / annulation",
          "UC-4  Médecins & disponibilités"]
droite = ["UC-5  Flux clinique du jour",
          "UC-6  Administration des cliniques",
          "UC-7  Tableaux de bord & statistiques"]
carte(s, Inches(0.45), Inches(2.2), Inches(6.1), Inches(3.0), "Les diagrammes", gauche, BLEU)
carte(s, Inches(6.75), Inches(2.2), Inches(6.1), Inches(3.0), " ", droite, BLEU)
carte(s, Inches(0.45), Inches(5.35), Inches(12.4), Inches(1.4),
      "Ce qu'ils montrent / lien avec le projet",
      ["Les acteurs, les fonctionnalités importantes et les interactions acteur-système (include / extend)",
       "Chaque diagramme = une fonctionnalité Must du cahier des charges ; ils orientent le découpage en Epics"],
      SARCELLE)
pied(s, "Souleymane", "1 min")
notes(s, "Passons au cœur du dossier de conception. Comme MediPlan est une application "
         "web, la consigne demandait des cas d'utilisation, un diagramme de classes et "
         "des diagrammes de séquence. Pour les cas d'utilisation, nous en avons produit 7, "
         "alors que le minimum était de 5. Ils couvrent tout le périmètre : la vue "
         "d'ensemble, l'authentification, la réservation, les disponibilités, le flux du "
         "jour, l'administration et les statistiques. Chaque diagramme montre les acteurs, "
         "les fonctionnalités et leurs interactions, et correspond à une fonctionnalité "
         "prioritaire du cahier des charges. Je laisse Zakaria présenter les classes et "
         "les séquences.")

# =========================================================================
# DIAPO 7 — Conception (2) classes & sequence  (ZAKARIA)
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 6, "Dossier de conception (2/2)", "Diagramme de classes & de séquence")
carte(s, Inches(0.45), Inches(1.7), Inches(6.1), Inches(5.05),
      "Diagramme de classes",
      ["Structure principale de la plateforme",
       "Classes : Utilisateur, Patient, Médecin, Administrateur,",
       "   Clinique, RendezVous, Disponibilite, Specialite, Notification",
       "Attributs, méthodes et relations entre classes",
       " ",
       "Lien : sert de base à la base de données PostgreSQL",
       "et au développement des fonctionnalités essentielles"], BLEU)
carte(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(5.05),
      "Diagramme de séquence",
      ["5 participants : Patient, Interface Web, API,",
       "   Base de données, Service de notification",
       "Connexion : vérification e-mail / mot de passe → JWT",
       "Réservation : vérifie que le créneau est libre",
       "   → enregistre, sinon message d'erreur (anti-doublon)",
       "Annulation : statut « annulé » (historique conservé)",
       "Notification après chaque opération"], SARCELLE)
pied(s, "Zakaria (classes) · Larbi (séquence)", "1 min 30")
notes(s, "[ZAKARIA] Merci. Je présente le diagramme de classes que j'ai réalisé. Il "
         "décrit la structure de la plateforme : les classes principales sont "
         "Utilisateur, Patient, Médecin, Administrateur, Clinique, RendezVous, "
         "Disponibilite, Specialite et Notification, avec leurs attributs, méthodes et "
         "relations. Il sert directement de base à la conception de notre base de données "
         "PostgreSQL. "
         "[LARBI] De mon côté, j'ai réalisé le diagramme de séquence, qui montre les "
         "échanges entre cinq participants : le patient, l'interface web, l'API, la base "
         "de données et le service de notification. On y voit la connexion qui retourne un "
         "jeton JWT, la réservation qui vérifie d'abord que le créneau est libre pour "
         "éviter les doubles réservations, et l'annulation qui change simplement le statut "
         "afin de garder l'historique.")

# =========================================================================
# DIAPO 8 — GitHub  (ZAKARIA)
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 7, "Dépôt GitHub", "Structure · emplacement des fichiers · accès")
rect(s, Inches(0.45), Inches(1.7), Inches(6.1), Inches(5.05), RGBColor(0x1E, 0x29, 0x3B))
tb, tf = zone_texte(s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(4.6))
arbo = ["mediplan/", "  README.md", "  docs/", "    cahier-des-charges/",
        "    consignes/", "    conception/        ← livrables",
        "      00-introduction.md", "      cas-utilisation/  (7 UC)",
        "      erd/   (MCD / ERD)", "      images/  (.mmd + .png)",
        "      README.md", "  (frontend/ & backend/ : phases 2-3)"]
premier = True
for ligne in arbo:
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    premier = False
    p.space_after = Pt(2)
    r = p.add_run(); r.text = ligne
    r.font.name = "Consolas"; r.font.size = Pt(13)
    surligne = ("conception" in ligne) or ("UC" in ligne) or ("ERD" in ligne)
    r.font.color.rgb = RGBColor(0x8B, 0xE9, 0xC8) if surligne else BLANC
carte(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(2.45),
      "Ce qui est déposé",
      ["Dossier docs/conception/ : tous les livrables",
       "7 cas d'utilisation + ERD (Mermaid + PNG)",
       "Explications écrites pour chaque diagramme",
       "README qui sert de table des matières"], SARCELLE)
carte(s, Inches(6.75), Inches(4.35), Inches(6.1), Inches(2.4),
      "Accès & remise (compte pour des points)",
      ["Dépôt privé partagé",
       "Professeure invitée comme collaboratrice (lecture)",
       "Lien GitHub + lien docs/conception/ soumis sur eCité"], BLEU)
pied(s, "Zakaria", "1 min")
# (Zakaria poursuit sur le dépôt GitHub)
notes(s, "Côté GitHub, voici comment notre dépôt est organisé. À la racine, un README et "
         "un dossier docs qui contient le cahier des charges, les consignes et surtout le "
         "dossier conception. C'est là que se trouvent tous nos livrables : les 7 cas "
         "d'utilisation, le diagramme entité-association, les fichiers Mermaid avec leurs "
         "images PNG, et un README qui sert de table des matières. Les dossiers frontend "
         "et backend seront ajoutés aux phases 2 et 3. Enfin, le dépôt est privé mais nous "
         "avons invité la professeure en lecture, et nous avons soumis sur eCité le lien "
         "du dépôt et le lien direct vers docs/conception.")

# =========================================================================
# DIAPO 9 — Travail realise / en cours / bloquant  (ZAKARIA)
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 8, "Travail réalisé, en cours et points à clarifier", None)
carte(s, Inches(0.45), Inches(1.7), Inches(4.0), Inches(5.05),
      "✅  Terminé",
      ["Cahier des charges",
       "7 diagrammes de cas d'utilisation",
       "Diagramme de classes",
       "Diagrammes de séquence",
       "Diagramme entité-association (ERD)",
       "Explications écrites",
       "Mise en place Jira & GitHub",
       "Soumission eCité"], SARCELLE)
carte(s, Inches(4.65), Inches(1.7), Inches(4.0), Inches(5.05),
      "⏳  En cours",
      ["Réintégration des diagrammes de",
       "   classes et de séquence dans",
       "   docs/conception/ (finalisés",
       "   dans un espace de travail",
       "   séparé)",
       "Export PNG final de ces deux",
       "   diagrammes dans le dépôt"], BLEU_CLAIR)
carte(s, Inches(8.85), Inches(1.7), Inches(4.0), Inches(5.05),
      "⚠  À clarifier / bloquant",
      ["Confirmer l'acceptation de",
       "   l'invitation GitHub par la",
       "   professeure",
       "Valider les règles de délai",
       "   d'annulation (ex. 24 h)",
       "Confirmer les sprints natifs",
       "   Jira (au-delà des labels)"], GRIS)
pied(s, "Larbi", "1 min")
notes(s, "Faisons le point honnêtement. Côté terminé : le cahier des charges, les 7 cas "
         "d'utilisation, le diagramme de classes, les séquences, l'ERD, toutes les "
         "explications écrites, ainsi que la mise en place de Jira, GitHub et la "
         "soumission eCité. En cours : les diagrammes de classes et de séquence ont été "
         "finalisés dans un espace de travail séparé ; il nous reste à les réintégrer "
         "proprement dans le dossier docs/conception du dépôt avec leurs images. Enfin, "
         "trois points à clarifier : confirmer que la professeure a bien accepté "
         "l'invitation, valider la règle exacte du délai d'annulation, et transformer nos "
         "labels de sprint en sprints natifs Jira.")

# =========================================================================
# DIAPO 10 — Prochaines etapes  (ZAKARIA)
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s)
bandeau_titre(s, 9, "Prochaines étapes", "Priorités pour la suite")
carte(s, Inches(0.45), Inches(1.7), Inches(6.1), Inches(2.4),
      "Finir de boucler le Sprint 1",
      ["Réintégrer classes & séquence dans le dépôt",
       "Vérifier l'accès de la professeure",
       "Passer les derniers tickets en « Terminé »"], SARCELLE)
carte(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(2.4),
      "Démarrer le Sprint 2 — Fondations",
      ["Squelettes Angular + NestJS",
       "Conteneurisation Docker Compose",
       "Authentification JWT + contrôle d'accès RBAC"], BLEU)
carte(s, Inches(0.45), Inches(4.25), Inches(12.4), Inches(2.5),
      "Ajustements & priorités",
      ["Prototyper tôt la gestion des disponibilités (risque technique identifié)",
       "Garder Jira à jour en continu : statuts, responsables, progression réelle",
       "Maintenir la relecture croisée en binôme pour la qualité des livrables",
       "Priorité n° 1 du projet : le cœur métier de réservation (anti-double-réservation)"],
      BLEU_CLAIR)
pied(s, "Larbi", "30 s")
notes(s, "Pour la suite : d'abord, finir de boucler le Sprint 1 en réintégrant les "
         "derniers diagrammes et en vérifiant l'accès de la professeure. Ensuite, nous "
         "démarrons le Sprint 2, consacré aux fondations techniques : mettre en place les "
         "squelettes Angular et NestJS, la conteneurisation Docker, et l'authentification "
         "avec JWT et RBAC. Côté ajustements, nous voulons prototyper tôt la gestion des "
         "disponibilités, qui est notre principal risque technique, garder Jira à jour en "
         "continu, et continuer à nous relire en binôme. Notre priorité numéro un reste le "
         "cœur métier : la réservation sans doublon.")

# =========================================================================
# DIAPO 11 — Conclusion / Questions
# =========================================================================
s = prs.slides.add_slide(BLANK)
fond(s, BLEU)
rect(s, 0, Inches(3.5), LARGEUR, Pt(4), SARCELLE)
tb, tf = zone_texte(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.4))
para(tf, "Merci de votre attention", taille=40, couleur=BLANC, gras=True,
     premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
tb, tf = zone_texte(s, Inches(1.0), Inches(3.8), Inches(11.3), Inches(1.6))
para(tf, "Questions ?", taille=28, couleur=BLEU_PALE, gras=True,
     premier=True, aligne=PP_ALIGN.CENTER, espace_avant=0)
para(tf, "Équipe MediPlan — Souleymane DIALLO  ·  Zakaria Lahouiri  ·  Larbi Saib",
     taille=16, couleur=BLEU_PALE2, aligne=PP_ALIGN.CENTER, espace_avant=14)
notes(s, "Voilà qui conclut notre présentation du Sprint 1. Pour résumer : la conception "
         "est terminée, notre projet est structuré dans Jira et GitHub, et nous sommes "
         "prêts à démarrer le développement au Sprint 2. Nous sommes maintenant "
         "disponibles tous les trois pour répondre à vos questions.")

# --- Sauvegarde ----------------------------------------------------------
sortie = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "MediPlan-Sprint1-Presentation.pptx")
prs.save(sortie)
print("Genere :", sortie)
print("Nombre de diapos :", len(prs.slides._sldIdLst))
